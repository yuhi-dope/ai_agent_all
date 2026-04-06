"""pptx_generator マイクロエージェント。提案書JSONからPPTXスライドを生成する。"""

import io
import logging
import time
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from workers.micro.models import MicroAgentInput, MicroAgentOutput, MicroAgentError

logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════════════════════
# デザインシステム
# ═══════════════════════════════════════════════════════════════
PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)        # ダークネイビー
PRIMARY_LIGHT = RGBColor(0x2D, 0x2D, 0x4A)  # ミディアムネイビー
ACCENT = RGBColor(0xE9, 0x45, 0x60)         # レッド
ACCENT_LIGHT = RGBColor(0xFD, 0xE8, 0xEB)   # 薄レッド
SUCCESS = RGBColor(0x00, 0x96, 0x5B)         # グリーン
SUCCESS_LIGHT = RGBColor(0xE6, 0xF7, 0xF0)  # 薄緑
WARNING = RGBColor(0xF0, 0x8C, 0x00)        # オレンジ
DANGER = RGBColor(0xDC, 0x35, 0x45)         # レッド
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x2D, 0x2D, 0x3F)
GRAY_800 = RGBColor(0x33, 0x37, 0x3D)
GRAY_600 = RGBColor(0x6C, 0x75, 0x7D)
GRAY_400 = RGBColor(0xAD, 0xB5, 0xBD)
GRAY_200 = RGBColor(0xE9, 0xEC, 0xEF)
GRAY_100 = RGBColor(0xF8, 0xF9, 0xFA)

FONT_NAME = "Meiryo"

# 16:9 ワイドスクリーン
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ═══════════════════════════════════════════════════════════════
# ヘルパー関数
# ═══════════════════════════════════════════════════════════════

def _set_bg(slide: Any, color: RGBColor) -> None:
    """スライドの背景色を設定する。"""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide: Any, l: int, t: int, w: int, h: int,
          fill: RGBColor | None = None, line: RGBColor | None = None) -> Any:
    """矩形シェイプを追加する。"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _rounded(slide: Any, l: int, t: int, w: int, h: int,
             fill: RGBColor | None = None, line: RGBColor | None = None) -> Any:
    """角丸矩形シェイプを追加する。"""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _tb(slide: Any, l: int, t: int, w: int, h: int,
        text: str, sz: int = 16, color: RGBColor = DARK,
        bold: bool = False, align: int = PP_ALIGN.LEFT) -> Any:
    """テキストボックスを追加する。"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    return box


def _add_p(tf: Any, text: str, sz: int = 14, color: RGBColor = DARK,
           bold: bool = False, align: int = PP_ALIGN.LEFT,
           space: int | None = None) -> Any:
    """テキストフレームに段落を追加する。"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.alignment = align
    if space is not None:
        p.space_before = Pt(space)
    return p


def _header_bar(slide: Any, title: str, subtitle: str | None = None) -> None:
    """統一ヘッダーバーを追加する。"""
    _rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), PRIMARY)
    _rect(slide, Inches(0), Inches(1.1), SLIDE_WIDTH, Pt(3), ACCENT)
    _tb(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
        title, sz=28, color=WHITE, bold=True)
    if subtitle:
        _tb(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.4),
            subtitle, sz=13, color=GRAY_400)


def _footer(slide: Any, num: int, total: int = 10) -> None:
    """統一フッターを追加する。"""
    _rect(slide, Inches(0), Inches(7.15), SLIDE_WIDTH, Pt(1), GRAY_200)
    _tb(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.3),
        "Confidential — シャチョツー（社長2号）ご提案書", sz=9, color=GRAY_400)
    _tb(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
        f"{num} / {total}", sz=9, color=GRAY_400, align=PP_ALIGN.RIGHT)


def _card(slide: Any, l: int, t: int, w: int, h: int,
          fill: RGBColor = WHITE, border: RGBColor = GRAY_200,
          shadow: bool = True) -> Any:
    """カード型コンテナを追加する。"""
    if shadow:
        _rect(slide, l + Pt(3), t + Pt(3), w, h, GRAY_200)
    return _rounded(slide, l, t, w, h, fill, border)


def _icon_num(slide: Any, l: int, t: int, num: int,
              color: RGBColor = ACCENT, size: float = 0.65) -> None:
    """数字付き丸アイコンを追加する。"""
    from pptx.enum.shapes import MSO_SHAPE
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(size), Inches(size))
    o.fill.solid()
    o.fill.fore_color.rgb = color
    o.line.fill.background()
    otf = o.text_frame
    otf.word_wrap = False
    op = otf.paragraphs[0]
    op.text = str(num)
    op.font.size = Pt(int(size * 24))
    op.font.color.rgb = WHITE
    op.font.bold = True
    op.font.name = FONT_NAME
    op.alignment = PP_ALIGN.CENTER


def _commafy(v: int | float | None) -> str:
    """数値をカンマ区切り文字列にする。"""
    if v is None:
        return "0"
    return f"{int(v):,}"


# ═══════════════════════════════════════════════════════════════
# スライド生成関数
# ═══════════════════════════════════════════════════════════════

def _slide_cover(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド1: 表紙"""
    cover = data.get("cover", {})
    company = cover.get("target_company", "御社")
    date = cover.get("date", "")
    title = cover.get("title", f"{company} 様向け ご提案書")
    subtitle = cover.get("subtitle", "AI社長秘書「シャチョツー」導入のご提案")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PRIMARY)

    # 上部アクセントライン
    _rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), ACCENT)

    # ドキュメント種別
    _tb(s, Inches(1.2), Inches(1.2), Inches(5), Inches(0.5),
        "ご提案書  |  シャチョツー", sz=15, color=GRAY_400)

    # メインタイトル
    _rect(s, Inches(1.2), Inches(2.2), Inches(3.5), Pt(3), ACCENT)
    _tb(s, Inches(1.2), Inches(2.5), Inches(10), Inches(1.2),
        title, sz=36, color=WHITE, bold=True)

    # サブタイトル
    _tb(s, Inches(1.2), Inches(4.0), Inches(10), Inches(0.8),
        subtitle, sz=20, color=RGBColor(0xE0, 0x80, 0x90))

    # メタ情報
    meta = f"提出先:  {company} 様\n提出元:  シャチョツー株式会社\n日付:    {date}"
    _tb(s, Inches(1.2), Inches(5.8), Inches(5), Inches(1.2),
        meta, sz=14, color=GRAY_400)

    _tb(s, Inches(8), Inches(6.3), Inches(4.5), Inches(0.5),
        "CONFIDENTIAL", sz=13, color=RGBColor(0xC5, 0x9D, 0x33),
        bold=True, align=PP_ALIGN.RIGHT)


def _slide_pain_points(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド2: 御社の課題認識"""
    pain_points = data.get("pain_points", [])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "御社の課題認識", "ヒアリングをもとに整理した現状の課題")

    # 課題を最大5件まで表示（優先度でソート）
    priority_order = {"high": 0, "medium": 1, "low": 2}
    sorted_pains = sorted(
        pain_points[:5],
        key=lambda x: priority_order.get(x.get("priority", "medium"), 1),
    )

    for i, pain in enumerate(sorted_pains):
        y = Inches(1.5 + i * 1.1)
        priority = pain.get("priority", "medium")
        priority_color = {
            "high": DANGER, "medium": WARNING, "low": GRAY_600
        }.get(priority, GRAY_600)

        # 優先度バッジ
        _rect(s, Inches(0.8), y, Pt(5), Inches(0.9), priority_color)

        # カテゴリ + 説明
        category = pain.get("category", "")
        description = pain.get("description", "")
        impact = pain.get("impact", "")

        _tb(s, Inches(1.1), y, Inches(3.5), Inches(0.4),
            category, sz=16, color=PRIMARY, bold=True)
        _tb(s, Inches(1.1), y + Inches(0.35), Inches(5.5), Inches(0.5),
            description, sz=13, color=GRAY_800)
        _tb(s, Inches(7.5), y + Inches(0.1), Inches(5), Inches(0.7),
            f"放置リスク: {impact}", sz=12, color=DANGER)

    _footer(s, 2)


def _slide_concept(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド3: シャチョツーとは（コンセプト図）"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "シャチョツーとは", "中小企業向け AI社長秘書サービス")

    # メインメッセージ
    msg = _card(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.9), ACCENT_LIGHT, ACCENT)
    tf = msg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "社長の頭の中を AI がデジタル化し、社員全員がいつでも参照できる「もう一人の社長」"
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 3つの機能カード
    features = [
        ("ブレイン", "社長の脳をデジタル化",
         "暗黙知・判断基準をAIが学習\n社員がQ&Aで即座に参照\n退職してもナレッジが残る"),
        ("BPO", "業界特化の業務自動化",
         "見積作成・書類生成を自動化\n業界ごとに最適化されたAI\nヒューマンチェック付きで安心"),
        ("デジタルツイン", "会社の状態を可視化",
         "5次元で会社の健康状態を把握\nリスクの早期発見・対策提案\nWhat-ifシミュレーション"),
    ]

    for i, (title, subtitle, desc) in enumerate(features):
        x = Inches(0.8 + i * 4.1)
        c = _card(s, x, Inches(2.8), Inches(3.8), Inches(4.0))

        _icon_num(s, x + Inches(0.25), Inches(2.95), i + 1, ACCENT)
        _tb(s, x + Inches(1.1), Inches(3.0), Inches(2.5), Inches(0.4),
            title, sz=20, color=PRIMARY, bold=True)
        _tb(s, x + Inches(0.3), Inches(3.5), Inches(3.3), Inches(0.4),
            subtitle, sz=14, color=GRAY_600)
        _rect(s, x + Inches(0.3), Inches(3.9), Inches(3.2), Pt(1), GRAY_200)
        _tb(s, x + Inches(0.3), Inches(4.1), Inches(3.3), Inches(2.2),
            desc, sz=13, color=GRAY_800)

    _footer(s, 3)


def _slide_brain(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド4: ブレイン機能"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "ブレイン機能", "社長の脳をデジタル化し、組織の知恵に変える")

    steps = [
        ("Step 1", "ナレッジ取り込み", "社長への音声インタビュー\n既存マニュアル・議事録の読み込み\nQ&A形式での追加学習"),
        ("Step 2", "AI構造化", "暗黙知を体系的に整理\n判断基準・意思決定ルールを抽出\n業界ノウハウとして構造化"),
        ("Step 3", "社員が活用", "チャットで社長に質問する感覚\n根拠つきの回答を即座に取得\n24時間いつでも利用可能"),
    ]

    for i, (step, title, desc) in enumerate(steps):
        x = Inches(0.8 + i * 4.1)
        c = _card(s, x, Inches(1.5), Inches(3.8), Inches(4.5))

        _tb(s, x + Inches(0.3), Inches(1.7), Inches(3.3), Inches(0.4),
            step, sz=13, color=ACCENT, bold=True)
        _tb(s, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(0.5),
            title, sz=20, color=PRIMARY, bold=True)
        _rect(s, x + Inches(0.3), Inches(2.6), Inches(3.2), Pt(1), GRAY_200)
        _tb(s, x + Inches(0.3), Inches(2.9), Inches(3.3), Inches(2.5),
            desc, sz=14, color=GRAY_800)

    # 矢印代わりの接続
    for i in range(2):
        x = Inches(4.4 + i * 4.1)
        _tb(s, x, Inches(3.5), Inches(0.8), Inches(0.6),
            ">>>", sz=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    _footer(s, 4)


def _slide_bpo(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド5: BPO機能"""
    solution_map = data.get("solution_map", [])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "BPO機能", "業種特化のAI業務自動化")

    if not solution_map:
        _tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
            "業種に最適化されたAIが、定型業務を自動化します。",
            sz=18, color=GRAY_800)
        _footer(s, 5)
        return

    # solution_mapから最大6件表示
    items = solution_map[:6]
    cols = min(3, len(items))
    rows = (len(items) + cols - 1) // cols

    for idx, sol in enumerate(items):
        col = idx % cols
        row = idx // cols
        x = Inches(0.8 + col * 4.1)
        y = Inches(1.5 + row * 2.8)

        c = _card(s, x, y, Inches(3.8), Inches(2.4))

        pain = sol.get("pain_point", "")
        solution = sol.get("solution", "")
        module = sol.get("module", "")
        effect = sol.get("effect", "")

        _tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.4),
            f"課題: {pain}", sz=12, color=DANGER)
        _rect(s, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Pt(1), GRAY_200)
        _tb(s, x + Inches(0.3), y + Inches(0.7), Inches(3.3), Inches(0.8),
            solution, sz=14, color=PRIMARY, bold=True)
        _tb(s, x + Inches(0.3), y + Inches(1.5), Inches(3.3), Inches(0.3),
            f"モジュール: {module}", sz=11, color=ACCENT)
        _tb(s, x + Inches(0.3), y + Inches(1.9), Inches(3.3), Inches(0.3),
            f"効果: {effect}", sz=11, color=SUCCESS)

    _footer(s, 5)


def _slide_modules(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド6: 導入モジュール提案"""
    modules = data.get("modules", [])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "導入モジュール提案", "御社に最適なモジュールの組み合わせ")

    if not modules:
        _tb(s, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
            "ヒアリング内容をもとに最適なモジュールをご提案します。",
            sz=18, color=GRAY_800)
        _footer(s, 6)
        return

    for i, mod in enumerate(modules[:4]):
        x = Inches(0.8 + (i % 2) * 6.2)
        y = Inches(1.5 + (i // 2) * 2.8)

        c = _card(s, x, y, Inches(5.8), Inches(2.4))

        name = mod.get("name", "")
        desc = mod.get("description", "")
        price = mod.get("monthly_price", 0)
        features = mod.get("key_features", [])

        _tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.5), Inches(0.5),
            name, sz=20, color=PRIMARY, bold=True)
        _tb(s, x + Inches(3.8), y + Inches(0.2), Inches(1.8), Inches(0.5),
            f"¥{_commafy(price)}/月", sz=16, color=ACCENT, bold=True,
            align=PP_ALIGN.RIGHT)

        _tb(s, x + Inches(0.3), y + Inches(0.7), Inches(5.2), Inches(0.5),
            desc, sz=13, color=GRAY_600)

        _rect(s, x + Inches(0.3), y + Inches(1.2), Inches(5.2), Pt(1), GRAY_200)

        features_text = "  |  ".join(features[:4])
        _tb(s, x + Inches(0.3), y + Inches(1.4), Inches(5.2), Inches(0.8),
            features_text, sz=12, color=GRAY_800)

    _footer(s, 6)


def _slide_pricing(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド7: 料金シミュレーション"""
    pricing = data.get("pricing", {})
    modules = data.get("modules", [])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "料金シミュレーション", "御社向けの料金プラン")

    # 料金テーブル
    col_w = [Inches(5.0), Inches(2.5), Inches(2.5)]
    headers = ["モジュール", "月額料金", "備考"]
    ht = Inches(0.55)
    y0 = Inches(1.5)

    # ヘッダー行
    for j, (h, cw) in enumerate(zip(headers, col_w)):
        x = Inches(1.5) + sum(col_w[:j])
        _rect(s, x, y0, cw, ht, PRIMARY)
        _tb(s, x, y0 + Inches(0.08), cw, Inches(0.4),
            h, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # データ行
    for i, mod in enumerate(modules[:6]):
        y = y0 + ht * (i + 1)
        bg = WHITE if i % 2 == 0 else GRAY_100
        name = mod.get("name", "")
        price = f"¥{_commafy(mod.get('monthly_price', 0))}"
        note = mod.get("description", "")[:30]

        vals = [name, price, note]
        for j, (v, cw) in enumerate(zip(vals, col_w)):
            x = Inches(1.5) + sum(col_w[:j])
            _rect(s, x, y, cw, ht, bg)
            lpad = Inches(0.3) if j == 0 else Inches(0)
            _tb(s, x + lpad, y + Inches(0.08), cw - lpad, Inches(0.4),
                v, sz=14,
                color=ACCENT if j == 1 else GRAY_800,
                bold=(j == 1),
                align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # 合計行
    total_y = y0 + ht * (min(len(modules), 6) + 1)
    total_w = sum(col_w)
    _rect(s, Inches(1.5), total_y, total_w, Inches(0.7), ACCENT_LIGHT)

    monthly = pricing.get("modules_total", 0)
    annual = pricing.get("annual_total", 0)
    discount = pricing.get("discount_note", "")

    _tb(s, Inches(1.8), total_y + Inches(0.1), Inches(4), Inches(0.5),
        "月額合計", sz=18, color=PRIMARY, bold=True)
    _tb(s, Inches(6.5), total_y + Inches(0.1), Inches(2.5), Inches(0.5),
        f"¥{_commafy(monthly)}/月", sz=20, color=ACCENT, bold=True,
        align=PP_ALIGN.CENTER)

    if discount:
        _tb(s, Inches(1.8), total_y + Inches(0.6), Inches(8), Inches(0.4),
            discount, sz=12, color=SUCCESS)

    # 年額表示
    _tb(s, Inches(1.8), total_y + Inches(1.0), Inches(8), Inches(0.4),
        f"年額換算: ¥{_commafy(annual)}", sz=14, color=GRAY_600)

    _footer(s, 7)


def _slide_roi(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド8: ROI試算（Before/After比較）"""
    roi = data.get("roi_estimate", {})
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "ROI試算", "投資対効果の定量分析")

    current_cost = roi.get("current_cost_monthly", 0)
    after_cost = roi.get("after_cost_monthly", 0)
    savings = roi.get("savings_monthly", 0)
    payback = roi.get("payback_months", 0)
    basis = roi.get("calculation_basis", "")

    # Before/After 比較カード
    # Before
    _card(s, Inches(0.8), Inches(1.5), Inches(5.0), Inches(2.5), WHITE, DANGER)
    _tb(s, Inches(1.1), Inches(1.7), Inches(4.5), Inches(0.5),
        "Before（現状）", sz=20, color=DANGER, bold=True)
    _tb(s, Inches(1.1), Inches(2.3), Inches(4.5), Inches(0.5),
        f"月間コスト: ¥{_commafy(current_cost)}", sz=24, color=DANGER, bold=True)
    _tb(s, Inches(1.1), Inches(3.0), Inches(4.5), Inches(0.5),
        "手作業中心・属人的な運用", sz=14, color=GRAY_600)

    # After
    _card(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5), WHITE, SUCCESS)
    _tb(s, Inches(7.3), Inches(1.7), Inches(5.0), Inches(0.5),
        "After（導入後）", sz=20, color=SUCCESS, bold=True)
    _tb(s, Inches(7.3), Inches(2.3), Inches(5.0), Inches(0.5),
        f"月間コスト: ¥{_commafy(after_cost)}", sz=24, color=SUCCESS, bold=True)
    _tb(s, Inches(7.3), Inches(3.0), Inches(5.0), Inches(0.5),
        "AI活用・標準化された運用", sz=14, color=GRAY_600)

    # 矢印（テキスト代替）
    _tb(s, Inches(5.5), Inches(2.3), Inches(1.5), Inches(0.5),
        ">>>", sz=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # 削減効果サマリー
    effect_box = _card(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.5),
                       ACCENT_LIGHT, ACCENT)
    tf = effect_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"月間削減額: ¥{_commafy(savings)}  |  投資回収期間: {payback}ヶ月"
    p.font.size = Pt(22)
    p.font.color.rgb = PRIMARY
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # 算出根拠
    if basis:
        _tb(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6),
            f"算出根拠: {basis}", sz=12, color=GRAY_600)

    _footer(s, 8)


def _slide_timeline(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド9: 導入スケジュール"""
    timeline = data.get("timeline", [])
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, WHITE)
    _header_bar(s, "導入スケジュール", "段階的な導入で確実に成果を出す")

    if not timeline:
        default_timeline = [
            {"phase": "Phase 1", "period": "1-2週目", "tasks": ["ヒアリング", "初期設定"], "milestone": "アカウント開設"},
            {"phase": "Phase 2", "period": "3-4週目", "tasks": ["ナレッジ取り込み", "テスト運用"], "milestone": "ブレイン稼働"},
            {"phase": "Phase 3", "period": "5-6週目", "tasks": ["BPO設定", "業務フロー最適化"], "milestone": "BPO稼働"},
            {"phase": "Phase 4", "period": "7-8週目", "tasks": ["本番運用", "効果測定"], "milestone": "本番移行完了"},
        ]
        timeline = default_timeline

    for i, phase in enumerate(timeline[:5]):
        x = Inches(0.8)
        y = Inches(1.5 + i * 1.15)

        phase_name = phase.get("phase", f"Phase {i + 1}")
        period = phase.get("period", "")
        tasks = phase.get("tasks", [])
        milestone = phase.get("milestone", "")

        # フェーズバー
        bar_color = ACCENT if i % 2 == 0 else PRIMARY_LIGHT
        _rect(s, x, y, Inches(2.0), Inches(0.9), bar_color)
        _tb(s, x + Inches(0.1), y + Inches(0.05), Inches(1.8), Inches(0.4),
            phase_name, sz=14, color=WHITE, bold=True)
        _tb(s, x + Inches(0.1), y + Inches(0.4), Inches(1.8), Inches(0.4),
            period, sz=12, color=RGBColor(0xCC, 0xCC, 0xCC))

        # タスク
        tasks_text = " / ".join(tasks[:3])
        _tb(s, Inches(3.0), y + Inches(0.1), Inches(6.0), Inches(0.7),
            tasks_text, sz=14, color=GRAY_800)

        # マイルストーン
        if milestone:
            _tb(s, Inches(9.5), y + Inches(0.1), Inches(3.5), Inches(0.7),
                f">>> {milestone}", sz=13, color=SUCCESS, bold=True)

    _footer(s, 9)


def _slide_next_steps(prs: Presentation, data: dict[str, Any]) -> None:
    """スライド10: Next Steps（CTA）"""
    cover = data.get("cover", {})
    company = cover.get("target_company", "御社")

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, PRIMARY)

    # 上部アクセントライン
    _rect(s, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), ACCENT)

    _tb(s, Inches(1.2), Inches(1.5), Inches(10), Inches(0.8),
        "Next Steps", sz=36, color=WHITE, bold=True)

    _rect(s, Inches(1.2), Inches(2.3), Inches(3.0), Pt(3), ACCENT)

    steps = [
        ("1", "無料トライアルのお申込み", "ブレイン機能を3ヶ月間 ¥30,000/月でお試し"),
        ("2", "初回ヒアリング（60分）", "社長の頭の中をAIに取り込む初回セッション"),
        ("3", "テスト運用開始", "2週間でQ&A機能が使える状態に"),
        ("4", "効果測定・本格導入判断", "数値で効果を確認してから正式契約"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.8 + i * 0.95)
        _icon_num(s, Inches(1.2), y, num, ACCENT)
        _tb(s, Inches(2.1), y, Inches(5), Inches(0.4),
            title, sz=18, color=WHITE, bold=True)
        _tb(s, Inches(2.1), y + Inches(0.4), Inches(8), Inches(0.4),
            desc, sz=14, color=GRAY_400)

    # CTA
    _tb(s, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6),
        f"{company} 様の「もう一人の社長」を、一緒につくりましょう。",
        sz=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# メインエントリポイント
# ═══════════════════════════════════════════════════════════════

def _build_pptx(proposal_data: dict[str, Any]) -> bytes:
    """提案書JSONからPPTXバイナリを生成する。

    Args:
        proposal_data: sales_proposal.py の出力形式に準拠したJSON辞書。

    Returns:
        PPTXファイルのバイナリデータ。
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _slide_cover(prs, proposal_data)
    _slide_pain_points(prs, proposal_data)
    _slide_concept(prs, proposal_data)
    _slide_brain(prs, proposal_data)
    _slide_bpo(prs, proposal_data)
    _slide_modules(prs, proposal_data)
    _slide_pricing(prs, proposal_data)
    _slide_roi(prs, proposal_data)
    _slide_timeline(prs, proposal_data)
    _slide_next_steps(prs, proposal_data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def run_pptx_generator(input: MicroAgentInput) -> MicroAgentOutput:
    """
    提案書JSONからPPTXスライドを生成する。

    payload:
        proposal_data (dict): sales_proposal.py の出力形式に準拠した提案書JSON。
            cover, pain_points, solution_map, modules, pricing, roi_estimate, timeline を含む。
        output_path (str, optional): 保存先ファイルパス。指定がなければメモリ内のみ。

    result:
        pptx_bytes (bytes): 生成されたPPTXバイナリ
        size_kb (float): ファイルサイズ (KB)
        slide_count (int): スライド枚数
        output_path (str | None): 保存先パス（指定があった場合）
    """
    start_ms = int(time.time() * 1000)
    agent_name = "pptx_generator"

    try:
        proposal_data: dict[str, Any] = input.payload.get("proposal_data", {})
        output_path: str | None = input.payload.get("output_path")

        if not proposal_data:
            raise MicroAgentError(agent_name, "input_validation", "proposal_data が空です")

        pptx_bytes = _build_pptx(proposal_data)
        size_kb = round(len(pptx_bytes) / 1024, 1)

        # ファイル保存（オプション）
        saved_path: str | None = None
        if output_path:
            from pathlib import Path
            # パストラバーサル防止
            real = Path(output_path).resolve()
            if ".." in str(real):
                raise MicroAgentError(agent_name, "security", "不正なファイルパスです")
            real.parent.mkdir(parents=True, exist_ok=True)
            real.write_bytes(pptx_bytes)
            saved_path = str(real)
            logger.info(f"PPTX saved to: {saved_path}")

        duration_ms = int(time.time() * 1000) - start_ms
        logger.info(f"PPTX generated: {size_kb} KB, {duration_ms}ms")

        return MicroAgentOutput(
            agent_name=agent_name,
            success=True,
            result={
                "pptx_bytes": pptx_bytes,
                "size_kb": size_kb,
                "slide_count": 10,
                "output_path": saved_path,
            },
            confidence=1.0,
            cost_yen=0.0,
            duration_ms=duration_ms,
        )

    except MicroAgentError:
        raise
    except Exception as e:
        duration_ms = int(time.time() * 1000) - start_ms
        logger.error(f"pptx_generator error: {e}")
        return MicroAgentOutput(
            agent_name=agent_name,
            success=False,
            result={"error": str(e)},
            confidence=0.0,
            cost_yen=0.0,
            duration_ms=duration_ms,
        )
