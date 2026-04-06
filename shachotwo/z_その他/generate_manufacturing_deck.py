"""製造業ドメイン知識大全 — プレゼンテーション v4
   全スライド: グリッド統一 / テキスト幅修正 / 高さ揃え / 値の短縮表記
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Tokens ──
BG     = RGBColor(0x0B, 0x11, 0x1E)
BG2    = RGBColor(0x07, 0x0B, 0x15)
CARD   = RGBColor(0x15, 0x1F, 0x32)
STRIPE = RGBColor(0x1B, 0x27, 0x3C)
BLUE   = RGBColor(0x3B, 0x82, 0xF6)
TEAL   = RGBColor(0x14, 0xB8, 0xA6)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
ROSE   = RGBColor(0xF4, 0x3F, 0x5E)
W      = RGBColor(0xFF, 0xFF, 0xFF)
L1     = RGBColor(0xE2, 0xE8, 0xF0)
L2     = RGBColor(0x94, 0xA3, 0xB8)
L3     = RGBColor(0x64, 0x74, 0x8B)

F  = "Hiragino Sans W3"
FB = "Hiragino Sans W6"

prs = Presentation()
prs.slide_width  = SW = Inches(13.333)
prs.slide_height = SH = Inches(7.5)

# ── Grid constants ──
ML = Inches(0.6)          # margin left
MR = Inches(0.6)          # margin right
COL1_W = Inches(5.8)      # single wide column
COL2_W = Inches(5.8)      # second column
GAP = Inches(0.35)        # gap between columns
COL1_L = ML               # col1 left
COL2_L = ML + COL1_W + GAP  # col2 left = 6.75"

HDR_T  = Inches(0.45)     # header top
BODY_T = Inches(1.6)      # body top (after header)
SRC_T  = Inches(7.05)     # source line top

# 3-column
C3W = Inches(3.85)
C3G = Inches(0.3)
C3_1 = ML
C3_2 = ML + C3W + C3G
C3_3 = ML + (C3W + C3G) * 2

# ── Primitives ──
def _bg(s, c=BG):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def _r(s, l, t, w, h, f, rad=0.03):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background()
    if rad: sh.adjustments[0] = rad
    return sh

def _ln(s, l, t, w, c=BLUE):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()

def _t(s, l, t, w, h, txt, sz=12, c=W, b=False, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = FB if b else F; p.alignment = al

def _ml(s, l, t, w, h, lines, sz=11, c=L1, sp=Pt(4)):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = F; p.space_after = sp

# ── Compound ──
def hdr(s, title, sub=None):
    _ln(s, ML, HDR_T, Inches(0.8))
    _t(s, ML, HDR_T + Inches(0.15), Inches(11), Inches(0.45), title, 26, W, True)
    if sub:
        _t(s, ML, HDR_T + Inches(0.6), Inches(11), Inches(0.3), sub, 12, L2)

def src(s, text):
    _t(s, ML, SRC_T, Inches(12), Inches(0.2), text, 8, L3)

def row(s, l, t, w, key, val, i=0, ksz=11, vsz=11):
    """Table row - key left, val right. FULL width for value."""
    f = STRIPE if i % 2 == 0 else CARD
    _r(s, l, t, w, Inches(0.3), f, 0.02)
    kw = w * 0.55
    vw = w * 0.42
    _t(s, l + Inches(0.12), t + Inches(0.03), kw, Inches(0.25), key, ksz, L1)
    _t(s, l + w * 0.55, t + Inches(0.03), vw, Inches(0.25), val, vsz, W, True, PP_ALIGN.RIGHT)

def card_hdr(s, l, t, w, title, clr=BLUE, sz=13):
    """Card header line + title. Returns y after title."""
    _ln(s, l + Inches(0.15), t + Inches(0.12), Inches(0.5), clr)
    _t(s, l + Inches(0.15), t + Inches(0.22), w - Inches(0.3), Inches(0.3), title, sz, W, True)
    return t + Inches(0.6)


# ═══════════════════════════════════════════════════════════════
# S1 — タイトル
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, BG2)
_r(s, Inches(0), Inches(3.15), SW, Pt(2), BLUE)
_t(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.8),
   "製造業ドメイン知識大全", 48, W, True, PP_ALIGN.CENTER)
_t(s, Inches(0.5), Inches(2.5), Inches(12.3), Inches(0.4),
   "中小製造業（従業員10-300名）の業務・文化・課題を網羅的に理解する", 17, L2, al=PP_ALIGN.CENTER)
for i, lb in enumerate(["全14章", "24業種", "見積・生産管理", "品質・設備", "財務・法規制"]):
    _r(s, Inches(1.7 + i * 2.1), Inches(3.8), Inches(1.8), Inches(0.36), CARD, 0.15)
    _t(s, Inches(1.7 + i * 2.1), Inches(3.82), Inches(1.8), Inches(0.32), lb, 10, L1, al=PP_ALIGN.CENTER)
_t(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.3),
   "対象: 営業 / PM / エンジニア / 経営者", 13, L2, al=PP_ALIGN.CENTER)
src(s, "出典: シャチョツー ナレッジベース g_02a_製造業ドメイン知識大全.md")

# ═══════════════════════════════════════════════════════════════
# S2 — 目次
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "目次")
chs = ["中小製造業とは何か","製造プロセスの全体像","見積の完全解説","生産管理の実態",
       "品質管理の実務","設備と保全","人と組織の課題","デジタル化状況",
       "業界トレンド","商習慣・取引慣行","物流・梱包","法規制",
       "資金繰り・財務","用語集"]
for i, ch in enumerate(chs):
    col, rw = i % 2, i // 2
    x, y = Inches(0.8 + col * 6.3), Inches(1.6 + rw * 0.68)
    clr = BLUE if col == 0 else TEAL
    _r(s, x, y, Inches(0.5), Inches(0.4), clr, 0.12)
    _t(s, x, y + Inches(0.04), Inches(0.5), Inches(0.32), f"{i+1}", 13, W, True, PP_ALIGN.CENTER)
    _t(s, x + Inches(0.65), y + Inches(0.04), Inches(5), Inches(0.32), ch, 15, L1)

# ═══════════════════════════════════════════════════════════════
# S3 — 業種分類 (3×2)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "中小製造業とは何か", "JSIC中分類24業種から中小が多い主要業種")

items = [
    ("金属製品製造業", "~43,000社", "切削・プレス・板金・鍛造・鋳造\n「丸物屋」「角物屋」で棲み分け", BLUE),
    ("生産用機械器具", "数千社", "工作機械・産業ロボット・半導体装置\n景気循環が激しい", TEAL),
    ("食料品製造業", "HACCP必須", "OEM製造が多い。配合展開型見積\nアレルゲン管理が大コスト", AMBER),
    ("電子部品", "多品種少量", "SMT実装・コネクタ・センサー\n1日数十回品種切替。RoHS/REACH", BLUE),
    ("輸送用機械", "IATF 16949", "自動車部品が圧倒的\n年次コストダウン 2-5%/年", TEAL),
    ("その他", "化学/繊維/ゴム", "バッチ生産・環境規制\n金型費100万-5,000万円", AMBER),
]
ROW_H = Inches(2.35)
for i, (nm, tag, desc, clr) in enumerate(items):
    col, rw = i % 3, i // 3
    x = Inches(0.5 + col * 4.15)
    y = BODY_T + Inches(rw * (ROW_H.inches + 0.25))
    _r(s, x, y, C3W, ROW_H, CARD)
    _ln(s, x + Inches(0.15), y + Inches(0.12), Inches(0.5), clr)
    _t(s, x + Inches(0.15), y + Inches(0.25), Inches(2.3), Inches(0.28), nm, 14, W, True)
    _r(s, x + Inches(2.6), y + Inches(0.22), Inches(1.0), Inches(0.26), clr, 0.12)
    _t(s, x + Inches(2.6), y + Inches(0.23), Inches(1.0), Inches(0.24), tag, 8, W, True, PP_ALIGN.CENTER)
    _ml(s, x + Inches(0.15), y + Inches(0.7), Inches(3.5), Inches(1.4), desc.split("\n"), 11, L1, Pt(5))
src(s, "出典: 総務省「日本標準産業分類」/ 経済産業省「工業統計」/ 中小企業白書")

# ═══════════════════════════════════════════════════════════════
# S4 — 財務構造
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "中小製造業の財務構造", "売上10億円・金属加工会社（30名）の典型的P/L")

# 左: P/L
_r(s, COL1_L, BODY_T, COL1_W, Inches(5.1), CARD)
_t(s, COL1_L + Inches(0.15), BODY_T + Inches(0.12), Inches(3), Inches(0.3), "典型的P/L構成", 13, W, True)
pl = [("売上高","10億円","100%",True),("材料費","3.5億円","35%",False),
      ("外注加工費","1.0億円","10%",False),("労務費","2.5億円","25%",False),
      ("製造経費","1.5億円","15%",False)]
for i,(n,v,p,bd) in enumerate(pl):
    y = BODY_T + Inches(0.55 + i * 0.28)
    c = W if bd else L1
    _t(s, COL1_L + Inches(0.3), y, Inches(2), Inches(0.25), n, 11, c, bd)
    _t(s, COL1_L + Inches(2.5), y, Inches(1.5), Inches(0.25), v, 11, c, bd, PP_ALIGN.RIGHT)
    _t(s, COL1_L + Inches(4.3), y, Inches(1), Inches(0.25), p, 10, L2, al=PP_ALIGN.RIGHT)

_r(s, COL1_L + Inches(0.3), BODY_T + Inches(1.98), Inches(5.2), Pt(1), L3)
res = [("粗利益","1.5億円","15%",False),("販管費","1.0億円","10%",False),("営業利益","5,000万円","5%",True)]
for i,(n,v,p,bd) in enumerate(res):
    y = BODY_T + Inches(2.1 + i * 0.28)
    c = BLUE if n == "営業利益" else (W if bd else L1)
    _t(s, COL1_L + Inches(0.3), y, Inches(2), Inches(0.25), n, 11, c, bd)
    _t(s, COL1_L + Inches(2.5), y, Inches(1.5), Inches(0.25), v, 11, c, bd, PP_ALIGN.RIGHT)
    _t(s, COL1_L + Inches(4.3), y, Inches(1), Inches(0.25), p, 10, c if bd else L2, al=PP_ALIGN.RIGHT)

# 左下: 利益率低い理由
_t(s, COL1_L + Inches(0.15), BODY_T + Inches(3.2), Inches(5), Inches(0.3), "なぜ利益率が低いのか", 12, W, True)
_ml(s, COL1_L + Inches(0.3), BODY_T + Inches(3.55), Inches(5.2), Inches(1.4), [
    "値引き圧力 — 年次コストダウン2-5%/年（10年で20-40%↓）",
    "材料高騰 — 転嫁に半年かかる。自社負担が常態化",
    "設備償却 — MC1台 2,000-5,000万円を10年で償却",
    "人件費の硬直性 — 受注減でも技能者は解雇不可",
], 10, L1, Pt(4))

# 右: 業種別粗利率
_r(s, COL2_L, BODY_T, COL2_W, Inches(2.7), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "業種別 粗利率", TEAL)
for i,(n,v) in enumerate([("板金加工","25-35%"),("電子部品実装","25-40%"),("金属切削","20-30%"),
                          ("射出成形","20-30%"),("プレス加工","15-25%"),("食品製造","15-25%")]):
    row(s, COL2_L + Inches(0.1), y0 + Inches(i * 0.32), COL2_W - Inches(0.2), n, v, i)

# 右下: 設備投資
_r(s, COL2_L, BODY_T + Inches(2.9), COL2_W, Inches(2.2), CARD)
card_hdr(s, COL2_L, BODY_T + Inches(2.9), COL2_W, "設備投資の目安", AMBER)
_ml(s, COL2_L + Inches(0.15), BODY_T + Inches(3.5), COL2_W - Inches(0.3), Inches(1.4), [
    "売上高の5-10%が投資目安",
    "MC1台: 2,000-5,000万円 / 5軸MC: 5,000万-1.5億",
    "中古設備: 新品の1/3-1/2で入手可能",
    "借入金: 月商の3-6ヶ月分が平均（金利0.5-2.0%）",
], 10, L1, Pt(4))

src(s, "出典: TKC経営指標 / 中小企業実態基本調査 / 中小企業白書")

# ═══════════════════════════════════════════════════════════════
# S5 — 組織構造と文化
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "組織構造と文化", "30名規模の金属加工会社 — 「1人=1部署」の世界")

# 左: 組織図
_r(s, COL1_L, BODY_T, COL1_W, Inches(5.1), CARD)
roles = [
    ("社長（1名）", "営業+見積最終判断+経営全般を兼務", BLUE),
    ("工場長（1名）", "辞めたら工場が止まる。生産計画は頭の中", ROSE),
    ("見積担当（1-3名）", "最も属人的。60代ベテラン1人で回すことも", AMBER),
    ("現場作業員（15-20名）", "正社員10+パート3+派遣2+実習生2", TEAL),
    ("事務員（1-2名）", "経理・総務・受発注を1-2人で全部やる", L3),
]
for i,(role,desc,clr) in enumerate(roles):
    y = BODY_T + Inches(0.15 + i * 0.95)
    _r(s, COL1_L + Inches(0.15), y, Inches(0.06), Inches(0.72), clr)
    _t(s, COL1_L + Inches(0.4), y + Inches(0.05), Inches(5), Inches(0.25), role, 12, W, True)
    _t(s, COL1_L + Inches(0.4), y + Inches(0.35), Inches(5), Inches(0.3), desc, 10, L2)

# 右: 文化（5つ）— 同じ高さ
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
card_hdr(s, COL2_L, BODY_T, COL2_W, "5つの文化的特徴", TEAL)
culture = [
    ("OJT文化", "マニュアルなし。「背中を見て覚えろ」が5-10年", BLUE),
    ("暗黙知の塊", "数百〜数千のコツがベテランの頭に。退職で消失", TEAL),
    ("設備への愛着", "「30年モノだけど精度が出る」合理的更新が阻まれる", AMBER),
    ("変化への慎重さ", "「今までこれでやってきた」が最強の反論", ROSE),
    ("紙・FAX・電話", "図面はFAX、見積はExcel→印刷→FAX", L3),
]
for i,(kw,desc,clr) in enumerate(culture):
    y = BODY_T + Inches(0.6 + i * 0.85)
    _r(s, COL2_L + Inches(0.1), y, COL2_W - Inches(0.2), Inches(0.72), STRIPE, 0.02)
    _r(s, COL2_L + Inches(0.1), y, Inches(0.05), Inches(0.72), clr)
    _t(s, COL2_L + Inches(0.3), y + Inches(0.08), Inches(2.5), Inches(0.25), kw, 11, W, True)
    _t(s, COL2_L + Inches(0.3), y + Inches(0.36), Inches(5.2), Inches(0.3), desc, 10, L2)

# ═══════════════════════════════════════════════════════════════
# S6 — 製造プロセス（3列テーブル — 値を短縮）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "製造プロセスの全体像")

# Big nums — 1行に収まるサイズ
for i,(num,lbl,clr) in enumerate([("80%","受注生産(MTO)",BLUE),("10%","見込み生産(MTS)",TEAL),("10%","受注組立(ATO)",AMBER)]):
    x = Inches(0.7 + i * 2.8)
    _t(s, x, BODY_T, Inches(2.5), Inches(0.45), num, 28, clr, True)
    _t(s, x, BODY_T + Inches(0.45), Inches(2.5), Inches(0.25), lbl, 10, L2)

# 3 cards — 値を短く
procs = [
    ("切削加工", [
        ("汎用旋盤", "¥3.5-5千/h"), ("CNC旋盤", "¥5-7千/h"), ("複合加工機", "¥8-12千/h"),
        ("3軸MC", "¥6-9千/h"), ("5軸MC", "¥1-1.5万/h"), ("門型MC", "¥1.5-2.5万/h"),
        ("ワイヤーカット", "¥8-12千/h"),
    ], BLUE),
    ("板金・プレス・成形", [
        ("レーザー切断", "SUS t1 ¥3-5/mm"), ("TIG溶接", "¥10-20/mm"), ("MIG/MAG", "¥5-10/mm"),
        ("プレス金型", "50万-3千万"), ("射出成形金型", "100万-5千万"),
        ("ABS樹脂", "¥250-400/kg"), ("PEEK樹脂", "¥1.5-2.2万/kg"),
    ], TEAL),
    ("表面処理・検査", [
        ("亜鉛メッキ", "¥100-500/個"), ("アルマイト", "¥100-500/個"), ("硬質クロム", "¥500-2千/個"),
        ("焼入れ", "¥300-1.5千/個"), ("窒化処理", "¥800-3千/個"),
        ("三次元測定機", "精度±0.001mm"), ("ノギス", "精度±0.02mm"),
    ], AMBER),
]
CARD_T = BODY_T + Inches(0.9)
CARD_H = Inches(4.6)
for col,(title,rows,clr) in enumerate(procs):
    x = Inches(0.5 + col * 4.15)
    _r(s, x, CARD_T, C3W, CARD_H, CARD)
    y0 = card_hdr(s, x, CARD_T, C3W, title, clr)
    for i,(k,v) in enumerate(rows):
        row(s, x + Inches(0.08), y0 + Inches(i * 0.33), C3W - Inches(0.16), k, v, i, 10, 10)
src(s, "出典: 各種加工メーカーヒアリング / NCジャパン / 業界標準値")

# ═══════════════════════════════════════════════════════════════
# S7 — 材料
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "材料の基礎知識", "kg単価は2024年時点の参考値")

# 左: 鉄鋼
_r(s, COL1_L, BODY_T, COL1_W, Inches(3.0), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "鉄鋼", BLUE)
for i,(mat,use,pr) in enumerate([
    ("SS400","一般構造物","¥100-130/kg"),("S45C","シャフト・ギア","¥130-170/kg"),
    ("SCM435","高強度ボルト","¥180-230/kg"),("SUS304","耐食部品","¥380-450/kg"),
    ("SUS316","海水・薬品","¥480-550/kg"),("SKD11","冷間金型","¥800-1,200/kg"),
    ("SKH51","切削工具","¥2-3千/kg"),
]):
    y = y0 + Inches(i * 0.3)
    f = STRIPE if i % 2 == 0 else CARD
    _r(s, COL1_L + Inches(0.08), y, COL1_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL1_L + Inches(0.2), y + Inches(0.02), Inches(1.0), Inches(0.24), mat, 10, BLUE, True)
    _t(s, COL1_L + Inches(1.3), y + Inches(0.02), Inches(1.8), Inches(0.24), use, 10, L1)
    _t(s, COL1_L + Inches(3.5), y + Inches(0.02), Inches(2.0), Inches(0.24), pr, 10, W, True, PP_ALIGN.RIGHT)

# 右: アルミ・銅
_r(s, COL2_L, BODY_T, COL2_W, Inches(3.0), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "アルミ・銅・真鍮", TEAL)
for i,(mat,use,pr) in enumerate([
    ("A5052","板金・筐体","¥450-550/kg"),("A6063","アルミフレーム","¥500-600/kg"),
    ("A7075","航空宇宙","¥750-900/kg"),
    ("C1100","純銅（電極）","¥1,100-1,400/kg"),("C3604","快削黄銅","¥1,100-1,300/kg"),
]):
    y = y0 + Inches(i * 0.3)
    f = STRIPE if i % 2 == 0 else CARD
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.02), Inches(1.0), Inches(0.24), mat, 10, TEAL, True)
    _t(s, COL2_L + Inches(1.3), y + Inches(0.02), Inches(1.8), Inches(0.24), use, 10, L1)
    _t(s, COL2_L + Inches(3.5), y + Inches(0.02), Inches(2.0), Inches(0.24), pr, 10, W, True, PP_ALIGN.RIGHT)

# 下段注意
_r(s, ML, BODY_T + Inches(3.2), Inches(11.8), Inches(0.85), CARD)
_r(s, ML, BODY_T + Inches(3.2), Inches(0.05), Inches(0.85), AMBER)
_ml(s, ML + Inches(0.2), BODY_T + Inches(3.3), Inches(11.4), Inches(0.65), [
    "SUS304の注意: 加工硬化で工具費がSS400の2-3倍。切削条件の選定がシビア",
    "買い方: 定尺材（丸棒4m, 板914×1829mm）/ 黒皮vs.ミガキ / ロット差20-30% / ミルシート（自動車・航空・医療で必須）",
], 10, L1, Pt(3))
src(s, "出典: 材料商社カタログ / JFEスチール・日本製鉄 / 各種メーカーヒアリング（2024年）")

# ═══════════════════════════════════════════════════════════════
# S8 — 見積フロー
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "見積の完全解説 — 6ステップ")

steps = [
    ("1","引合い受領","図面+仕様が届く\n自社対応可否を判断",BLUE),
    ("2","図面読み取り","★最も属人的\n三面図→立体化\n暗黙要求を判断",AMBER),
    ("3","工程推定","★2番目に属人的\n形状→設備選定\n工程順序の決定",BLUE),
    ("4","工数算出","段取り+CT×数量\n旋盤30-45分\nMC45-60分",TEAL),
    ("5","コスト計算","材料+加工+外注\n+管理費+利益\n得意先別ルール",BLUE),
    ("6","提出","Excel→FAX\n社長の最終チェック",TEAL),
]
STEP_W = Inches(1.85)
STEP_GAP = Inches(0.25)
for i,(num,title,desc,clr) in enumerate(steps):
    x = Inches(0.3 + i * (STEP_W.inches + STEP_GAP.inches))
    _r(s, x + Inches(0.68), BODY_T, Inches(0.5), Inches(0.5), clr, 0.5)
    _t(s, x + Inches(0.68), BODY_T + Inches(0.08), Inches(0.5), Inches(0.38), num, 18, W, True, PP_ALIGN.CENTER)
    _t(s, x, BODY_T + Inches(0.6), STEP_W, Inches(0.25), title, 11, W, True, PP_ALIGN.CENTER)
    _r(s, x, BODY_T + Inches(0.9), STEP_W, Inches(2.2), CARD)
    _ml(s, x + Inches(0.1), BODY_T + Inches(1.0), Inches(1.65), Inches(2.0), desc.split("\n"), 10, L1, Pt(3))
    if i < 5:
        _t(s, x + STEP_W, BODY_T + Inches(0.12), Inches(0.25), Inches(0.3), "→", 11, L3, al=PP_ALIGN.CENTER)

# 下部ポイント
_r(s, Inches(0.3), BODY_T + Inches(3.4), Inches(12.7), Inches(1.0), CARD)
_ln(s, Inches(0.5), BODY_T + Inches(3.52), Inches(0.5), AMBER)
_t(s, Inches(0.5), BODY_T + Inches(3.6), Inches(12), Inches(0.25), "見積の核心 — 「図面に書いてないこと」の判断", 12, W, True)
_ml(s, Inches(0.5), BODY_T + Inches(3.9), Inches(12), Inches(0.4), [
    "バリ取り指示なし→でも残すとクレーム / 面取り指示なし→C0.2-0.3は常識 / メッキ厚を考慮した前寸法 / 得意先ごとの暗黙ルール",
], 10, L2, Pt(2))

# ═══════════════════════════════════════════════════════════════
# S9 — チャージレート
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "チャージレートと計算例")

# 左上: 計算式
_r(s, COL1_L, BODY_T, COL1_W, Inches(0.85), CARD)
_ml(s, COL1_L + Inches(0.15), BODY_T + Inches(0.1), COL1_W - Inches(0.3), Inches(0.7), [
    "チャージレート＝（設備コスト+労務費+間接費）÷ 稼働時間",
    "例: CNC旋盤3千万 → (償却300万+電力60万+人件費500万)÷2,000h ≒ ¥4,300/h",
], 9, L1, Pt(3))

# 左下: 計算例
_r(s, COL1_L, BODY_T + Inches(0.95), COL1_W, Inches(3.5), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(0.95), COL1_W, "計算例: シャフト φ30×100mm, S45C", TEAL)
for i,(proc,cond,time) in enumerate([
    ("外径荒加工","V=150, f=0.25, 2パス","0.58分"),("外径仕上げ","V=200, f=0.1","0.47分"),
    ("端面加工 ×2","","0.60分"),("面取り ×4","","0.40分"),("ワーク脱着","","0.50分"),
]):
    y = y0 + Inches(i * 0.32)
    f = STRIPE if i % 2 == 0 else CARD
    _r(s, COL1_L + Inches(0.08), y, COL1_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL1_L + Inches(0.2), y + Inches(0.02), Inches(1.5), Inches(0.24), proc, 10, L1)
    _t(s, COL1_L + Inches(1.8), y + Inches(0.02), Inches(2.0), Inches(0.24), cond, 9, L3)
    _t(s, COL1_L + Inches(4.0), y + Inches(0.02), Inches(1.5), Inches(0.24), time, 10, W, True, PP_ALIGN.RIGHT)
# 合計行
ty = y0 + Inches(5 * 0.32)
_r(s, COL1_L + Inches(0.08), ty, COL1_W - Inches(0.16), Inches(0.3), BLUE, 0.02)
_t(s, COL1_L + Inches(0.2), ty + Inches(0.02), Inches(3), Inches(0.25), "合計サイクルタイム", 11, W, True)
_t(s, COL1_L + Inches(4.0), ty + Inches(0.02), Inches(1.5), Inches(0.25), "≒ 2.1分/個", 12, W, True, PP_ALIGN.RIGHT)

# 右: 表面粗さ
_r(s, COL2_L, BODY_T, COL2_W, Inches(2.8), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "表面粗さとコスト", AMBER)
for i,(sym,ra,method,cost,clr) in enumerate([
    ("1","Ra 25-6.3","旋盤のまま","基準",L3),("2","Ra 3.2-1.6","仕上げ加工","+10-20%",TEAL),
    ("3","Ra 0.8-0.4","研磨","+30-50%",AMBER),("4","Ra 0.2-0.1","ラッピング","+100-300%",ROSE),
]):
    y = y0 + Inches(i * 0.42)
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.38), STRIPE if i%2==0 else CARD, 0.02)
    _r(s, COL2_L + Inches(0.08), y, Inches(0.05), Inches(0.38), clr)
    # 番号バッジ
    _r(s, COL2_L + Inches(0.2), y + Inches(0.06), Inches(0.3), Inches(0.26), clr, 0.12)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.07), Inches(0.3), Inches(0.24), sym, 10, W, True, PP_ALIGN.CENTER)
    _t(s, COL2_L + Inches(0.6), y + Inches(0.05), Inches(1.2), Inches(0.25), ra, 10, L1)
    _t(s, COL2_L + Inches(1.9), y + Inches(0.05), Inches(1.3), Inches(0.25), method, 10, L2)
    _t(s, COL2_L + Inches(3.5), y + Inches(0.05), Inches(2.0), Inches(0.25), cost, 10, W, True, PP_ALIGN.RIGHT)

# 右下: 段取り
_r(s, COL2_L, BODY_T + Inches(3.0), COL2_W, Inches(1.45), CARD)
y0 = card_hdr(s, COL2_L, BODY_T + Inches(3.0), COL2_W, "段取り時間の目安", BLUE)
for i,(m,t) in enumerate([("CNC旋盤","30-45分"),("マシニングセンタ","45-60分"),("5軸MC","60-90分")]):
    row(s, COL2_L + Inches(0.08), y0 + Inches(i * 0.3), COL2_W - Inches(0.16), m, t, i, 10, 10)

src(s, "出典: NCジャパン見積エンジン / 各種メーカー実績値 / JIS B 0601")

# ═══════════════════════════════════════════════════════════════
# S10 — 生産管理
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "生産管理の実態", "「工場長の頭の中」が最も普及した生産管理システム")

_r(s, COL1_L, BODY_T, COL1_W, Inches(2.7), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "生産計画の実態", BLUE)
_ml(s, COL1_L + Inches(0.15), y0, COL1_W - Inches(0.3), Inches(2.0), [
    "MC-1号機: 3/27 A社50個 → 3/28 B社20個 → 4/1 C社5個",
    "",
    "工場長が休むと「今日何を作るか」誰もわからない",
    "新しい注文の納期回答も工場長しかできない",
    "→ 中小製造業の最大の経営リスクの一つ",
], 10, L1, Pt(3))

_r(s, COL1_L, BODY_T + Inches(2.9), COL1_W, Inches(2.2), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(2.9), COL1_W, "山崩し（キャパ超過時の5対策）", TEAL)
_ml(s, COL1_L + Inches(0.15), y0, COL1_W - Inches(0.3), Inches(1.5), [
    "① 前倒し（空き日に移す）  ② 外注（マージン10-15%）",
    "③ 残業（割増25-50%, 月45h上限）  ④ 休日出勤（割増35%）",
    "⑤ 納期交渉（最終手段）",
], 10, L1, Pt(3))

# 右: 赤字受注
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "赤字受注の5つのメカニズム", ROSE)
for i,(cause,desc) in enumerate([
    ("見積ミス","サイクルタイム過小評価・段取り見落とし"),
    ("仕様変更","「やっぱりSUS304で」→ 追加請求できない"),
    ("不良発生","10%不良→材料+加工時間が丸損+追加生産"),
    ("材料高騰","見積時と発注時で鋼材が10-20%上昇"),
    ("値引き圧力","「5%引けば発注する」→ 利益率割れ"),
]):
    y = y0 + Inches(i * 0.82)
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.7), STRIPE if i%2==0 else CARD, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.06), Inches(2), Inches(0.22), f"❌ {cause}", 11, ROSE, True)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.34), Inches(5.2), Inches(0.3), desc, 10, L2)
src(s, "出典: 中小企業庁「下請取引改善」/ 中小企業白書")

# ═══════════════════════════════════════════════════════════════
# S11 — 品質管理
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "品質管理の実務")

# 左上: ISO
_r(s, COL1_L, BODY_T, COL1_W, Inches(2.4), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "ISO 9001「ISO疲れ」の実態", BLUE)
_ml(s, COL1_L + Inches(0.15), y0, COL1_W - Inches(0.3), Inches(1.6), [
    "取得率: 30名以上の製造業で約40%。動機は「得意先の要求」",
    "品質管理担当1人で全文書管理。年1回更新で1ヶ月",
    "審査費用: 年間50-100万円",
    "「書類のための書類」感が強い",
], 10, L1, Pt(4))

# 左下: SPC
_r(s, COL1_L, BODY_T + Inches(2.6), COL1_W, Inches(2.5), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(2.6), COL1_W, "SPC — 工程能力指数 Cpk", TEAL)
for i,(val,desc,clr) in enumerate([
    ("< 1.0","能力不足（全数検査）",ROSE),("≧ 1.33","良好（自動車最低要求）",AMBER),
    ("≧ 1.67","優良（航空・医療）",TEAL),("≧ 2.0","卓越（6σ）",BLUE),
]):
    y = y0 + Inches(i * 0.35)
    _r(s, COL1_L + Inches(0.08), y, COL1_W - Inches(0.16), Inches(0.3), STRIPE if i%2==0 else CARD, 0.02)
    _t(s, COL1_L + Inches(0.2), y + Inches(0.03), Inches(1.0), Inches(0.24), val, 10, clr, True)
    _t(s, COL1_L + Inches(1.4), y + Inches(0.03), Inches(4.0), Inches(0.24), desc, 10, L1)

# 右: なぜなぜ分析
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "不良対応 — 8D + なぜなぜ分析", AMBER)
_ml(s, COL2_L + Inches(0.15), y0, COL2_W - Inches(0.3), Inches(0.6), [
    "8D: 準備→チーム→問題記述→暫定対策→根本原因→恒久対策→検証→再発防止",
], 10, L2, Pt(2))

_t(s, COL2_L + Inches(0.15), y0 + Inches(0.7), Inches(5), Inches(0.25), "なぜなぜ分析の実例", 11, W, True)
for i,(q,a) in enumerate([
    ("なぜ外径が太い？","仕上げの切込み不足"),("なぜ切込み不足？","荒加工の寸法バラつき"),
    ("なぜバラつく？","チャック把持力が不安定"),("なぜ不安定？","爪が摩耗していた"),
    ("なぜ気づかない？","定期交換ルールなし"),
]):
    y = y0 + Inches(1.05 + i * 0.38)
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.33), STRIPE if i%2==0 else CARD, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.03), Inches(2.5), Inches(0.25), q, 10, AMBER)
    _t(s, COL2_L + Inches(2.7), y + Inches(0.03), Inches(2.8), Inches(0.25), f"→ {a}", 10, L1)

_r(s, COL2_L + Inches(0.08), y0 + Inches(2.98), COL2_W - Inches(0.16), Inches(0.3), BLUE, 0.02)
_t(s, COL2_L + Inches(0.2), y0 + Inches(3.0), COL2_W - Inches(0.4), Inches(0.25),
   "真因: 爪交換基準未設定 → 対策: 月1回測定, 0.05mm超で交換", 10, W, True)
src(s, "出典: ISO 9001:2015 / IATF 16949 / JIS Q 9001")

# ═══════════════════════════════════════════════════════════════
# S12 — 人と設備
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "人と組織の課題 / 設備と保全")

# 4 big nums
for i,(num,lbl,clr) in enumerate([
    ("1.5-2.0x","有効求人倍率",ROSE),("30-40%","3年以内離職率",AMBER),
    ("~30%","後継者未定",BLUE),("50代","現場平均年齢",TEAL)]):
    _t(s, Inches(0.6 + i * 3.1), BODY_T, Inches(2.8), Inches(0.4), num, 24, clr, True)
    _t(s, Inches(0.6 + i * 3.1), BODY_T + Inches(0.4), Inches(2.8), Inches(0.2), lbl, 9, L2)

# 左: 課題
_r(s, COL1_L, BODY_T + Inches(0.8), COL1_W, Inches(3.9), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(0.8), COL1_W, "人的課題", AMBER)
for i,(t,d) in enumerate([
    ("技能伝承の危機","ベテラン1人に数百〜数千のコツ。動画では伝わらない"),
    ("外国人実習生","3年で一人前→5年で帰国の構造問題"),
    ("後継者問題","社長の頭が最大の無形資産。M&A年3-4千件"),
    ("賃金格差","製造業380万 vs 全産業430万（約12%低い）"),
]):
    y = y0 + Inches(i * 0.68)
    _t(s, COL1_L + Inches(0.2), y, Inches(2.5), Inches(0.22), t, 10, W, True)
    _t(s, COL1_L + Inches(0.2), y + Inches(0.25), Inches(5.3), Inches(0.3), d, 9, L2)

# 右: 設備テーブル
_r(s, COL2_L, BODY_T + Inches(0.8), COL2_W, Inches(3.9), CARD)
y0 = card_hdr(s, COL2_L, BODY_T + Inches(0.8), COL2_W, "設備価格帯（新品 / 中古）", TEAL)
for i,(nm,nw,us) in enumerate([
    ("CNC旋盤","1,500-4,000万","500-2,000万"),("複合加工機","3,000-8,000万","1,500-4,000万"),
    ("3軸MC","2,000-5,000万","800-2,500万"),("5軸MC","5,000万-1.5億","2,500-8,000万"),
    ("門型MC","8,000万-3億","3,000万-1億"),("レーザー","3,000-8,000万","1,500-4,000万"),
    ("射出成形(350t)","2,000-4,000万","800-2,000万"),
]):
    y = y0 + Inches(i * 0.3)
    f = STRIPE if i%2==0 else CARD
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.02), Inches(1.8), Inches(0.24), nm, 10, L1)
    _t(s, COL2_L + Inches(2.2), y + Inches(0.02), Inches(1.5), Inches(0.24), nw, 10, W, b=True, al=PP_ALIGN.RIGHT)
    _t(s, COL2_L + Inches(4.0), y + Inches(0.02), Inches(1.5), Inches(0.24), us, 10, TEAL, al=PP_ALIGN.RIGHT)
src(s, "出典: 厚労省「職業安定業務統計」/ 中小企業白書 / 各社カタログ（2024年）")

# ═══════════════════════════════════════════════════════════════
# S13 — デジタル化 + トレンド
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "デジタル化状況 / 業界トレンド")

# 左上: IT
_r(s, COL1_L, BODY_T, COL1_W, Inches(2.4), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "IT導入率", BLUE)
for i,(nm,pct) in enumerate([("会計ソフト",90),("CAD/CAM",55),("生産管理",30),("見積ソフト",12),("MES",7)]):
    y = y0 + Inches(i * 0.3)
    _t(s, COL1_L + Inches(0.15), y, Inches(1.3), Inches(0.25), nm, 10, L1)
    _r(s, COL1_L + Inches(1.6), y + Inches(0.06), Inches(3.0), Inches(0.13), STRIPE, 0.2)
    _r(s, COL1_L + Inches(1.6), y + Inches(0.06), Inches(3.0 * pct / 100), Inches(0.13), BLUE, 0.2)
    _t(s, COL1_L + Inches(4.7), y, Inches(0.8), Inches(0.25), f"{pct}%", 10, W, True, PP_ALIGN.RIGHT)

# 左下: DX壁
_r(s, COL1_L, BODY_T + Inches(2.6), COL1_W, Inches(2.5), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(2.6), COL1_W, "DXの4つの壁", AMBER)
for i,(w,n) in enumerate([
    ("「今のやり方で回っている」","最も根深い"),("IT人材の不在","詳しい若手が兼任"),
    ("投資余力","利益率3-5%に5千万は重い"),("ロックイン警戒","過去の苦い経験"),
]):
    y = y0 + Inches(i * 0.35)
    _t(s, COL1_L + Inches(0.15), y, Inches(3.5), Inches(0.28), f"❌ {w}", 10, AMBER)
    _t(s, COL1_L + Inches(3.7), y, Inches(1.8), Inches(0.28), n, 9, L3)

# 右: トレンド
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "業界を揺るがす6つのトレンド", TEAL)
for i,(nm,desc,clr) in enumerate([
    ("多品種少量化","「100種類×各10個」が増加。段取り・見積負荷増",BLUE),
    ("短納期化","「2週間→1週間→3日」の圧縮",BLUE),
    ("材料価格高騰","SS400がコロナ前の1.3-1.5倍。全額転嫁わずか30%",AMBER),
    ("カーボンニュートラル","Scope 3算定要求がサプライヤーに波及",TEAL),
    ("国内回帰","TSMC熊本に象徴。中小に追い風だが人手不足",TEAL),
    ("EV化","エンジン部品3万点→EV1万点。2/3が不要に",ROSE),
]):
    y = y0 + Inches(i * 0.65)
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.55), STRIPE if i%2==0 else CARD, 0.02)
    _r(s, COL2_L + Inches(0.08), y, Inches(0.05), Inches(0.55), clr)
    _t(s, COL2_L + Inches(0.25), y + Inches(0.05), Inches(2.5), Inches(0.22), nm, 10, W, True)
    _t(s, COL2_L + Inches(0.25), y + Inches(0.28), Inches(5.2), Inches(0.22), desc, 9, L2)
src(s, "出典: 中小企業庁DX調査 / 経産省「素形材ビジョン」/ 日本鉄鋼連盟")

# ═══════════════════════════════════════════════════════════════
# S14 — 商習慣（3列）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "商習慣・取引慣行", "技術力だけでは仕事は来ない")

CH = Inches(5.1)
# 左
_r(s, C3_1, BODY_T, C3W, CH, CARD)
y0 = card_hdr(s, C3_1, BODY_T, C3W, "支払条件", BLUE)
for i,(term,note,clr) in enumerate([
    ("月末翌月末","最短30-60日",TEAL),("手形60日","良心的",TEAL),
    ("手形90日","最も一般的",AMBER),("手形120日","下請法ギリギリ",ROSE),
]):
    y = y0 + Inches(i * 0.35)
    _r(s, C3_1 + Inches(0.08), y, C3W - Inches(0.16), Inches(0.3), STRIPE if i%2==0 else CARD, 0.02)
    _r(s, C3_1 + Inches(0.08), y, Inches(0.04), Inches(0.3), clr)
    _t(s, C3_1 + Inches(0.2), y + Inches(0.03), Inches(1.5), Inches(0.24), term, 10, L1)
    _t(s, C3_1 + Inches(2.0), y + Inches(0.03), Inches(1.5), Inches(0.24), note, 10, clr, True, PP_ALIGN.RIGHT)
_ml(s, C3_1 + Inches(0.15), y0 + Inches(1.6), C3W - Inches(0.3), Inches(1.5), [
    "ファクタリング:",
    "  2社間: 手数料10-20%（非通知）",
    "  3社間: 手数料2-9%（承諾要）",
    "「手形より3-5%値引きして現金即払い」も",
], 9, L1, Pt(3))

# 中
_r(s, C3_2, BODY_T, C3W, CH, CARD)
y0 = card_hdr(s, C3_2, BODY_T, C3W, "値決めの力学", AMBER)
_ml(s, C3_2 + Inches(0.15), y0, C3W - Inches(0.3), Inches(2.0), [
    "相見積（3社競合）",
    "  本命1社+当て馬2社。価格横流し",
    "指値発注",
    "  「¥1,500/個で」根拠なし",
    "コストテーブル開示",
    "  原価丸裸→「ここを削れ」",
], 10, L1, Pt(3))
_t(s, C3_2 + Inches(0.15), y0 + Inches(2.2), C3W - Inches(0.3), Inches(0.25), "年次コストダウン", 11, W, True)
_ml(s, C3_2 + Inches(0.15), y0 + Inches(2.5), C3W - Inches(0.3), Inches(1.0), [
    "自動車: 毎年2-5%。5年で累計14%",
    "VA/VE提案→削減効果は親事業者が吸収",
], 9, L1, Pt(3))

# 右
_r(s, C3_3, BODY_T, C3W, CH, CARD)
y0 = card_hdr(s, C3_3, BODY_T, C3W, "書類の流れ", TEAL)
for i,doc in enumerate(["① 見積依頼(RFQ)","② 見積書","③ 注文書(PO)","④ 注文請書",
                        "⑤ 納品書","⑥ 検収","⑦ 請求書"]):
    y = y0 + Inches(i * 0.35)
    _r(s, C3_3 + Inches(0.08), y, C3W - Inches(0.16), Inches(0.3), STRIPE if i%2==0 else CARD, 0.02)
    _t(s, C3_3 + Inches(0.2), y + Inches(0.03), Inches(3), Inches(0.24), doc, 10, L1)
_r(s, C3_3 + Inches(0.08), y0 + Inches(2.65), C3W - Inches(0.16), Inches(0.35), CARD, 0.02)
_r(s, C3_3 + Inches(0.08), y0 + Inches(2.65), Inches(0.04), Inches(0.35), ROSE)
_t(s, C3_3 + Inches(0.2), y0 + Inches(2.68), Inches(3), Inches(0.28), "口頭発注は下請法3条違反", 10, ROSE, True)
src(s, "出典: 公正取引委員会「下請法運用基準」/ 中小企業庁")

# ═══════════════════════════════════════════════════════════════
# S15a — 物流・梱包（輸送+梱包の詳細）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "物流・梱包の暗黙ルール（1/2）", "見積書に載らないが、知らないと全数返品になることも")

# 左: 輸送形態
_r(s, COL1_L, BODY_T, COL1_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "輸送形態の使い分け", BLUE)
transport = [
    ("チャーター便（専用車）",
     "大型部品や精密部品の輸送。運送会社に1台まるごと手配。\n片道¥15,000-50,000。長尺物（3m超）はチャーター必須。\n「明日の朝イチで届けてくれ」という急ぎ案件で頻繁に使用"),
    ("路線便（佐川・西濃・福山）",
     "最も一般的。翌日-翌々日着。60サイズ¥800-160サイズ¥2,500。\n重量制限30kg/個。西濃運輸は重量物に強く製造業利用率が高い。\n毎日集荷に来る契約を結ぶのが普通"),
    ("混載便",
     "複数荷主の荷物をまとめて輸送。パレット単位が基本。\n¥3,000-15,000/パレット。路線便より安いが到着日にブレ（1-3日）"),
    ("自社便",
     "社長や営業マンがハイエース/軽トラで届ける。近距離（30分圏内）\nの急ぎ時。人件費を考えると最も割高だが「ライン停止」には勝てない"),
    ("赤帽（軽貨物チャーター）",
     "個人事業主の協同組合。当日配送可能。¥10,000-30,000。\n積載350kgまで。「さっき加工が終わった。今日中に検査に出したい」"),
]
for i,(name,desc) in enumerate(transport):
    y = y0 + Inches(i * 0.85)
    _t(s, COL1_L + Inches(0.15), y, Inches(5.5), Inches(0.22), name, 10, W, True)
    _ml(s, COL1_L + Inches(0.15), y + Inches(0.24), Inches(5.5), Inches(0.55), desc.split("\n"), 8, L2, Pt(1))

# 右: 梱包ルール
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "梱包の暗黙ルール — 知らないと全数返品", TEAL)
packing = [
    ("精密部品は1個ずつ個別包装",
     "面粗度Ra1.6以下の仕上げ面や嵌合部のある部品は、部品同士の\n接触で傷がつくためビニール袋+気泡緩衝材で個別包装が鉄則"),
    ("SUSと鉄は絶対に混ぜない",
     "異種金属を同じ箱に入れると「もらい錆」が発生。鉄の錆が\nステンレスに移り赤茶色のシミに。除去には酸洗いが必要。即クレーム"),
    ("メッキ後の部品は素手厳禁",
     "手の脂（指紋）が経時で変色やシミの原因に。ニトリル手袋必須。\n特にニッケルメッキ・クロムメッキは指紋が目立ちやすい"),
    ("アルミ部品は紙を巻く",
     "アルミは段ボールの酸性成分（リグニン）と反応して白錆（白い\n粉状の酸化物）が出る。中性紙やポリ袋で包んでから段ボールへ"),
    ("数量が多い場合は小分け",
     "100個なら10個ずつポリ袋に小分け+段ボール。検収担当に喜ばれる。\n見積書には載らない「暗黙のサービス」"),
    ("重量物は木枠梱包",
     "50kg超は木枠を組んでパレットにボルト止め。費用¥3,000-15,000。\n輸出品はISPM 15マーク（くん蒸処理済み）パレット必須"),
]
for i,(name,desc) in enumerate(packing):
    y = y0 + Inches(i * 0.7)
    _t(s, COL2_L + Inches(0.15), y, Inches(5.5), Inches(0.22), name, 10, W, True)
    _ml(s, COL2_L + Inches(0.15), y + Inches(0.24), Inches(5.5), Inches(0.45), desc.split("\n"), 8, L2, Pt(1))

# ═══════════════════════════════════════════════════════════════
# S15b — 納品書類・在庫管理の詳細
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "物流・梱包の暗黙ルール（2/2）", "納品書類・送り状・在庫管理の実態")

# 左: 納品書類
_r(s, COL1_L, BODY_T, COL1_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "納品伝票・送り状", BLUE)
docs_detail = [
    ("納品書のフォーマット",
     "品名・図番（図面番号）・数量・ロット番号・納品日を記載。\n得意先指定フォーマットがある場合はそれに従う。\n自社フォーマットでも最低限「図番」と「数量」は必須"),
    ("ミルシート・検査成績書の同梱",
     "ミルシート（材質証明書）: 鋼材メーカーが発行する公式書類。\n自動車・航空・医療では必須。得意先要求時に納品物に同梱。\n検査成績書: 自社での測定結果（寸法・硬度・表面粗さ等）"),
    ("送り状の作成",
     "「元払い」（発送側負担）が一般的。「着払い」や「第三者払い」\n（発注元が運賃負担、届け先は別工場）もある。\n第三者払いは自動車業界の支給品輸送でよく使われる"),
]
for i,(name,desc) in enumerate(docs_detail):
    y = y0 + Inches(i * 1.25)
    _t(s, COL1_L + Inches(0.15), y, Inches(5.5), Inches(0.22), name, 10, W, True)
    _ml(s, COL1_L + Inches(0.15), y + Inches(0.25), Inches(5.5), Inches(0.9), desc.split("\n"), 9, L2, Pt(2))

# 右: 在庫管理
_r(s, COL2_L, BODY_T, COL2_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "在庫管理の実態", AMBER)
inventory = [
    ("材料在庫（定尺材のストック）",
     "よく使うSS400丸棒（φ20-80）やSUS304フラットバー等を常備。\n月商の0.5-1ヶ月分が目安。特殊鋼（SKD11, SCM435等）は\nリードタイム2-4週間のため多めにストックする傾向"),
    ("仕掛品（加工途中の部品）",
     "工程間で停滞しがちなのが中小製造業の宿痾。\n「旋盤が終わったがフライスが空かない」「外注の熱処理待ち」等。\n仕掛品の金額把握が甘いと利益が見えなくなる"),
    ("端材管理",
     "「いつか使う」と取っておいた端材が棚を圧迫。φ50×L150mmを\n3年保管して結局使わないケースは珍しくない。\n「1年使わなかった端材は処分」等のルールが必要だが実行は困難"),
    ("完成品在庫・棚卸し",
     "得意先の内示（正式発注ではない口頭情報）で先行生産→内示が\n外れると不良在庫に。棚卸しは年1-2回。帳簿と実在庫の差異\n（棚差）の原因は出庫記録漏れ・端材未記録・不良廃棄未計上"),
]
for i,(name,desc) in enumerate(inventory):
    y = y0 + Inches(i * 1.0)
    _t(s, COL2_L + Inches(0.15), y, Inches(5.5), Inches(0.22), name, 10, W, True)
    _ml(s, COL2_L + Inches(0.15), y + Inches(0.25), Inches(5.5), Inches(0.7), desc.split("\n"), 9, L2, Pt(2))

src(s, "出典: 各種運送会社料金表 / 製造業ヒアリング / ISPM 15国際基準")

# ═══════════════════════════════════════════════════════════════
# S16 — 法規制
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "法規制・コンプライアンス")

_t(s, ML, BODY_T, Inches(3), Inches(0.4), "~8,000件/年", 24, ROSE, True)
_t(s, ML, BODY_T + Inches(0.4), Inches(5), Inches(0.2), "公正取引委員会の下請法指導件数", 10, L2)

# 左: 下請法+労安法
_r(s, COL1_L, BODY_T + Inches(0.8), COL1_W, Inches(2.2), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(0.8), COL1_W, "下請法の要点", BLUE)
_ml(s, COL1_L + Inches(0.15), y0, COL1_W - Inches(0.3), Inches(1.5), [
    "適用: 資本金3億超 → 3億以下への製造委託",
    "60日ルール: 受領日から60日以内に支払い（遅延利息14.6%）",
    "禁止: ① 受領拒否 ② 代金減額 ③ 返品 ④ 買いたたき",
    "3条書面: 発注時の書面交付義務（口頭発注は違法）",
], 10, L1, Pt(4))

_r(s, COL1_L, BODY_T + Inches(3.2), COL1_W, Inches(1.9), CARD)
y0 = card_hdr(s, COL1_L, BODY_T + Inches(3.2), COL1_W, "労働安全衛生法", TEAL)
_ml(s, COL1_L + Inches(0.15), y0, COL1_W - Inches(0.3), Inches(1.2), [
    "50名以上: 安全管理者・衛生管理者・産業医が必須",
    "特別教育: クレーン・フォーク・溶接・有機溶剤・粉じん",
    "特殊健康診断(年2回) / 2024年: ばく露低減措置の義務化",
], 10, L1, Pt(3))

# 右: 環境
_r(s, COL2_L, BODY_T + Inches(0.8), COL2_W, Inches(4.3), CARD)
y0 = card_hdr(s, COL2_L, BODY_T + Inches(0.8), COL2_W, "環境規制・CSR調達", AMBER)
for i,(nm,desc) in enumerate([
    ("RoHS指令","鉛・水銀・カドミウム等6物質の使用制限"),
    ("REACH規則","EU化学物質規制。SVHC含有報告義務"),
    ("化管法","PRTR制度。化学物質の排出・移動量を届出"),
    ("消防法","危険物の貯蔵・取扱い。指定数量以上は許可制"),
    ("CSR調達","人権・紛争鉱物・情報セキュリティの遵守要求"),
]):
    y = y0 + Inches(i * 0.62)
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.52), STRIPE if i%2==0 else CARD, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.04), Inches(1.5), Inches(0.22), nm, 10, AMBER, True)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.28), Inches(5.2), Inches(0.22), desc, 9, L2)
src(s, "出典: 公正取引委員会年次報告 / 厚労省 / 経産省化学物質規制")

# ═══════════════════════════════════════════════════════════════
# S17a — 資金繰り（手形・ファクタリング詳細 + 設備投資）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "資金繰り・財務の深掘り（1/2）", "倒産理由第2位「資金繰り難」の構造")

# 左: 手形・ファクタリング詳細
_r(s, COL1_L, BODY_T, COL1_W, Inches(5.1), CARD)
y0 = card_hdr(s, COL1_L, BODY_T, COL1_W, "手形・ファクタリングの実務", BLUE)

sections = [
    ("手形サイト（支払期日までの日数）", [
        ("60日", "比較的良心的。大手でも支払条件が良い取引先", TEAL),
        ("90日", "最も一般的。中小製造業の「標準」", AMBER),
        ("120日", "下請法ギリギリ。実質的に無利子借入を強いる構造", ROSE),
    ]),
    ("手形割引（銀行で早期現金化）", [
        ("割引率", "年利1.5-3.5%。信用力が高い振出人ほど低率", L1),
        ("計算例", "額面100万・90日・2.0% → 割引料¥4,932 → 手取り約99.5万", L1),
        ("注意", "不渡り時は買戻し義務あり（リスクは自社が負う）", ROSE),
    ]),
    ("ファクタリング（売掛債権の売却）", [
        ("2社間", "手数料10-20%。取引先に知られない。高額", AMBER),
        ("3社間", "手数料2-9%。取引先の承諾必要。安い", TEAL),
        ("注意", "恒常的に頼ると利益が飛ぶ。毎月利用は黄色信号", ROSE),
    ]),
]
y = y0
for sec_title, items in sections:
    _t(s, COL1_L + Inches(0.15), y, Inches(5.5), Inches(0.22), sec_title, 10, W, True)
    y += Inches(0.25)
    for label, desc, clr in items:
        _r(s, COL1_L + Inches(0.1), y, COL1_W - Inches(0.2), Inches(0.26), STRIPE, 0.02)
        _t(s, COL1_L + Inches(0.2), y + Inches(0.02), Inches(0.8), Inches(0.22), label, 9, clr, True)
        _t(s, COL1_L + Inches(1.1), y + Inches(0.02), Inches(4.4), Inches(0.22), desc, 9, L2)
        y += Inches(0.28)
    y += Inches(0.12)

# 右上: 設備投資
_r(s, COL2_L, BODY_T, COL2_W, Inches(2.4), CARD)
y0 = card_hdr(s, COL2_L, BODY_T, COL2_W, "設備投資の判断基準", TEAL)
for i,(term,result,clr) in enumerate([("3年以内回収","Go",TEAL),("3-5年回収","慎重検討",AMBER),("5年超","基本No",ROSE)]):
    y = y0 + Inches(i * 0.32)
    f = STRIPE if i%2==0 else CARD
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.02), Inches(2.5), Inches(0.24), term, 10, L1)
    _t(s, COL2_L + Inches(3.5), y + Inches(0.02), Inches(2.0), Inches(0.24), f"→ {result}", 10, clr, True, PP_ALIGN.RIGHT)
_ml(s, COL2_L + Inches(0.15), y0 + Inches(1.1), COL2_W - Inches(0.3), Inches(0.5), [
    "リース: 初期ゼロ・月額固定。技術変化が早い分野向き",
    "買取: 長期的にコスト安。5年以上使う見込みがあれば",
], 9, L2, Pt(2))

# 右下: 借入金（値幅を十分に確保）
_r(s, COL2_L, BODY_T + Inches(2.6), COL2_W, Inches(2.5), CARD)
y0 = card_hdr(s, COL2_L, BODY_T + Inches(2.6), COL2_W, "借入金の種類と使い分け", AMBER)
loans = [
    ("プロパー融資", "金利 0.5-1.5%", "優良企業のみ"),
    ("保証協会付", "金利 1.5-4.0%", "中小の主力調達手段"),
    ("日本政策金融公庫", "金利 0.5-2.0%", "創業融資に強い"),
    ("セーフティネット", "一般+別枠", "最大5.6億まで保証拡大"),
]
for i,(nm,rate,note) in enumerate(loans):
    y = y0 + Inches(i * 0.3)
    f = STRIPE if i%2==0 else CARD
    _r(s, COL2_L + Inches(0.08), y, COL2_W - Inches(0.16), Inches(0.28), f, 0.02)
    _t(s, COL2_L + Inches(0.2), y + Inches(0.02), Inches(1.8), Inches(0.24), nm, 9, L1)
    _t(s, COL2_L + Inches(2.1), y + Inches(0.02), Inches(1.5), Inches(0.24), rate, 9, W, True)
    _t(s, COL2_L + Inches(3.7), y + Inches(0.02), Inches(1.8), Inches(0.24), note, 8, L3)
_t(s, COL2_L + Inches(0.15), y0 + Inches(1.35), COL2_W - Inches(0.3), Inches(0.22),
   "適正借入: 月商の3倍以内が健全。6倍超は危険水域", 9, L2)

src(s, "出典: 全国銀行協会 / 日本政策金融公庫 / 中小企業庁 / 中小企業白書")

# ═══════════════════════════════════════════════════════════════
# S17b — 経営指標
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s)
hdr(s, "資金繰り・財務の深掘り（2/2）", "経営指標の読み方 — 銀行がどこを見ているか")

# 5つの指標を2列で表示
metrics = [
    ("自己資本比率", "純資産÷総資産×100",
     [("30%以上","健全。銀行が安心して貸せるライン",TEAL),
      ("10%以下","危険。赤字で債務超過に転落するリスク",ROSE),
      ("中小平均","35-40%（TKC経営指標）",L3)]),
    ("借入金月商倍率", "有利子負債÷月商",
     [("3倍以内","健全。月商5千万なら借入1.5億以内",TEAL),
      ("6倍超","危険。利息だけで月50-100万の負担",ROSE)]),
    ("流動比率", "流動資産÷流動負債×100",
     [("150%以上","健全。200%あれば余裕あり",TEAL),
      ("100%以下","危険。短期的な資金ショートリスク",ROSE)]),
    ("労働分配率", "人件費÷粗利益×100",
     [("50-60%","健全。人にも会社にも適切に分配",TEAL),
      ("70%超","危険。設備投資や内部留保の余裕なし",ROSE)]),
    ("債務償還年数", "(借入-現預金)÷(税引後利益+減価償却費)",
     [("10年以内","健全",TEAL)]),
]

for i,(name,formula,vals) in enumerate(metrics):
    col = i % 2
    rw = i // 2
    x = COL1_L if col == 0 else COL2_L
    w = COL1_W if col == 0 else COL2_W
    y_base = BODY_T + Inches(rw * 1.7)

    _r(s, x, y_base, w, Inches(1.55), CARD)
    _t(s, x + Inches(0.15), y_base + Inches(0.1), Inches(3), Inches(0.25), name, 12, W, True)
    _t(s, x + Inches(3.2), y_base + Inches(0.12), Inches(2.4), Inches(0.22), formula, 8, L3)

    for j,(lbl,desc,clr) in enumerate(vals):
        vy = y_base + Inches(0.45 + j * 0.3)
        _r(s, x + Inches(0.08), vy, w - Inches(0.16), Inches(0.28), STRIPE, 0.02)
        _t(s, x + Inches(0.2), vy + Inches(0.02), Inches(1.2), Inches(0.24), lbl, 9, clr, True)
        _t(s, x + Inches(1.5), vy + Inches(0.02), Inches(4.0), Inches(0.24), desc, 9, L2)

src(s, "出典: TKC経営指標 / 中小企業実態基本調査 / 全国銀行協会")

# ═══════════════════════════════════════════════════════════════
# S18 — まとめ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, BG2)
_ln(s, ML, Inches(0.4), Inches(0.8))
_t(s, ML, Inches(0.55), Inches(11), Inches(0.45), "まとめ — 製造業を理解するために", 26, W, True)

for i,(num,title,desc,clr) in enumerate([
    ("1","中小製造業の本質","「1人=1部署」「暗黙知の塊」「紙・FAX・電話依存」の世界。大企業の常識は通用しない。",BLUE),
    ("2","DXの最大機会","見積と生産管理が最も属人的で退職リスクが最大。ここにデジタル化の最大の機会。",TEAL),
    ("3","10年以内の臨界点","人手不足・技能伝承・後継者問題は待ったなし。現場平均年齢50代。",AMBER),
    ("4","取引の力学を理解せよ","支払条件・値決め・コストダウン・下請法。知らずに現場に刺さる提案はできない。",ROSE),
]):
    y = Inches(1.4 + i * 1.3)
    _r(s, ML, y, Inches(12.1), Inches(1.1), CARD)
    _r(s, ML, y, Inches(0.05), Inches(1.1), clr)
    _r(s, ML + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.5), clr, 0.12)
    _t(s, ML + Inches(0.25), y + Inches(0.3), Inches(0.5), Inches(0.4), num, 16, W, True, PP_ALIGN.CENTER)
    _t(s, ML + Inches(0.9), y + Inches(0.15), Inches(4), Inches(0.3), title, 16, W, True)
    _t(s, ML + Inches(0.9), y + Inches(0.5), Inches(10.8), Inches(0.45), desc, 11, L1)

_r(s, Inches(0), Inches(6.6), SW, Inches(0.9), RGBColor(0x05, 0x08, 0x12))
_t(s, Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.45),
   "この知識を武器に、製造業の真のパートナーになる。", 20, BLUE, True, PP_ALIGN.CENTER)

# ── Save ──
out = os.path.expanduser("~/code/ai_agent/製造業ドメイン知識大全.pptx")
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
