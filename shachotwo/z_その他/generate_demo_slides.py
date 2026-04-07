"""
シャチョツー 営業デモ資料 生成スクリプト
出力: shachotwo/c_事業計画/シャチョツー_営業デモ資料.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os

# =========================================
# カラー定義
# =========================================
NAVY = RGBColor(0x0F, 0x34, 0x60)
DARK_NAVY = RGBColor(0x16, 0x21, 0x3E)
RED = RGBColor(0xE9, 0x45, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)
LIGHT_NAVY = RGBColor(0x0A, 0x52, 0x94)

# フォント定義
# Google Slides互換優先: "Noto Sans JP" → "Hiragino Sans" → "Yu Gothic" → "Arial"
# pptxをGoogle Slidesで開く場合、Noto Sans JPが最も確実に表示される
import sys as _sys

def get_font():
    import platform
    if platform.system() == "Darwin":
        # macOS: Hiragino Sans（ローカルでの高品質表示）
        return "Hiragino Sans"
    return "Arial"

def get_google_slides_font():
    """Google Slides互換フォント（pptxをGドライブにアップロード時）"""
    return "Noto Sans JP"

FONT_NAME = get_font()
# Google Slides向けエクスポート時は FONT_NAME = get_google_slides_font() に切り替え
print(f"Using font: {FONT_NAME}")

# =========================================
# スライドサイズ（16:9 ワイドスクリーン）
# =========================================
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# =========================================
# ヘルパー関数
# =========================================

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, align=PP_ALIGN.LEFT,
                font_name=None):
    """テキストボックスを追加する"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name or FONT_NAME
    if color:
        run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """塗りつぶし矩形を追加する"""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_background(slide, color):
    """スライド背景色を設定する"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide, num, total=15, text_color=None):
    """スライド番号（右下）とフッター（左下）を追加する"""
    color = text_color or MID_GRAY
    # フッター「シャチョツー」
    add_textbox(
        slide,
        left=Cm(1), top=SLIDE_HEIGHT - Cm(1.2),
        width=Cm(8), height=Cm(0.9),
        text="シャチョツー",
        font_size=9, color=color, align=PP_ALIGN.LEFT
    )
    # スライド番号
    add_textbox(
        slide,
        left=SLIDE_WIDTH - Cm(3), top=SLIDE_HEIGHT - Cm(1.2),
        width=Cm(2.5), height=Cm(0.9),
        text=f"{num} / {total}",
        font_size=9, color=color, align=PP_ALIGN.RIGHT
    )


def add_divider(slide, top, color=RED):
    """水平区切り線を追加する"""
    from pptx.util import Pt as PtUnit
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Cm(1.5), top,
        SLIDE_WIDTH - Cm(1.5), top
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.5)


def add_section_title(slide, title, subtitle=None, bg_navy=False):
    """セクションタイトルバーを追加する"""
    color = WHITE if bg_navy else NAVY
    add_rect(slide, Cm(1.5), Cm(0.8), Cm(0.3), Cm(1.2), RED)
    add_textbox(
        slide,
        left=Cm(2.2), top=Cm(0.7),
        width=SLIDE_WIDTH - Cm(4), height=Cm(1.4),
        text=title,
        font_size=24, bold=True, color=color
    )
    if subtitle:
        add_textbox(
            slide,
            left=Cm(2.2), top=Cm(2.2),
            width=SLIDE_WIDTH - Cm(4), height=Cm(0.9),
            text=subtitle,
            font_size=13, color=MID_GRAY if not bg_navy else WHITE
        )


# =========================================
# スライド生成関数
# =========================================

def slide_01_cover(prs):
    """Slide 1: 表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_background(slide, NAVY)

    # アクセントライン（上部）
    add_rect(slide, 0, 0, SLIDE_WIDTH, Cm(0.6), RED)

    # メインキャッチコピー
    add_textbox(
        slide,
        left=Cm(3), top=Cm(3.5),
        width=SLIDE_WIDTH - Cm(6), height=Cm(3),
        text="月3万で社内に\nデジタル事務員を雇う",
        font_size=38, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # アンダーライン
    add_rect(slide, Cm(10), Cm(7.0), Cm(14), Cm(0.12), RED)

    # サブタイトル
    add_textbox(
        slide,
        left=Cm(3), top=Cm(7.5),
        width=SLIDE_WIDTH - Cm(6), height=Cm(1.5),
        text="シャチョツー（社長2号）",
        font_size=22, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF),
        align=PP_ALIGN.CENTER
    )

    # 説明文
    add_textbox(
        slide,
        left=Cm(3), top=Cm(9.2),
        width=SLIDE_WIDTH - Cm(6), height=Cm(1.2),
        text="会社の知恵をAIに移す。1,000社の経営者の知恵が使える。",
        font_size=14, color=RGBColor(0xAA, 0xBB, 0xDD),
        align=PP_ALIGN.CENTER
    )

    # 年号
    add_textbox(
        slide,
        left=Cm(3), top=SLIDE_HEIGHT - Cm(3.5),
        width=SLIDE_WIDTH - Cm(6), height=Cm(0.9),
        text="2026",
        font_size=13, color=MID_GRAY,
        align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, 1, text_color=RGBColor(0x66, 0x77, 0x99))


def slide_02_problem(prs):
    """Slide 2: こんなことありませんか？"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "こんなことありませんか？")
    add_divider(slide, Cm(2.3))

    problems = [
        ("社長への依存・属人化",
         "「社長に聞かないとわからない」が多すぎる。\n社長が休む・辞めると判断が止まる。毎日LINEが47件届く。"),
        ("ベテランが辞めたらノウハウが消える",
         "熟練社員の暗黙知は口伝え。退職すると経験が消える。\n後継者への技術承継・判断基準の移転ができていない。"),
        ("同じ質問に何度も答えている",
         "新人・若手が毎回同じことを聞いてくる。\n社長の時間が月20〜30時間、繰り返しの回答に消えている。"),
    ]

    icons = ["●", "●", "●"]
    top_positions = [Cm(3.2), Cm(7.5), Cm(11.8)]

    for i, (title, desc) in enumerate(problems):
        top = top_positions[i]
        # 背景ボックス
        add_rect(slide, Cm(1.5), top, SLIDE_WIDTH - Cm(3), Cm(3.8),
                 LIGHT_GRAY)

        # アイコン（●）
        add_textbox(
            slide, Cm(2.0), top + Cm(0.5),
            Cm(1.2), Cm(1.2),
            icons[i], font_size=22, bold=True, color=RED
        )
        # タイトル
        add_textbox(
            slide, Cm(3.5), top + Cm(0.4),
            SLIDE_WIDTH - Cm(5.5), Cm(1.0),
            title, font_size=18, bold=True, color=NAVY
        )
        # 説明
        add_textbox(
            slide, Cm(3.5), top + Cm(1.4),
            SLIDE_WIDTH - Cm(5.5), Cm(2.0),
            desc, font_size=12, color=DARK_NAVY
        )

    add_slide_number(slide, 2)


def slide_03_what(prs):
    """Slide 3: シャチョツーとは"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "シャチョツーとは",
                      subtitle="会社のデジタルツイン — 社長の知恵を丸ごとAIに移す")
    add_divider(slide, Cm(2.3))

    features = [
        ("🧠  ブレイン",
         "社長の判断基準・ノウハウをAIが学習。\n社員がAIに聞けば、社長の代わりに即回答。\nQ&A・テンプレート・書類作成が初日から使える。",
         RED),
        ("⚙️  BPO自動化",
         "見積・請求・契約・勤怠など業種特化の\n繰り返し業務をAIが自動処理。\n使うほど自動化レベルが段階的に上がる。",
         NAVY),
        ("📡  能動提案",
         "AIが先にリスクを検知して社長に通知。\n「来月の人員不足を検知しました」\n「受注過多でキャッシュフローが悪化します」",
         RGBColor(0x11, 0x88, 0x55)),
    ]

    positions = [Cm(1.5), Cm(12.3), Cm(23.1)]
    box_width = Cm(10.0)

    for i, (title, desc, color) in enumerate(features):
        left = positions[i]
        top = Cm(3.2)
        add_rect(slide, left, top, box_width, Cm(10.5), LIGHT_GRAY)
        add_rect(slide, left, top, box_width, Cm(0.25), color)

        add_textbox(
            slide, left + Cm(0.5), top + Cm(0.5),
            box_width - Cm(1), Cm(1.4),
            title, font_size=17, bold=True, color=NAVY
        )
        add_textbox(
            slide, left + Cm(0.5), top + Cm(2.0),
            box_width - Cm(1), Cm(7.5),
            desc, font_size=12, color=DARK_NAVY
        )

    add_slide_number(slide, 3)


def slide_04_network(prs):
    """Slide 4: データネットワーク効果"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Cm(0.4), RED)

    add_textbox(
        slide, Cm(3), Cm(1.5),
        SLIDE_WIDTH - Cm(6), Cm(2.0),
        "1,000社の知恵が使える",
        font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Cm(3), Cm(3.5),
        SLIDE_WIDTH - Cm(6), Cm(1.0),
        "データネットワーク効果 — 使うほど、全員が賢くなる",
        font_size=15, color=RGBColor(0xAA, 0xBB, 0xDD),
        align=PP_ALIGN.CENTER
    )

    points = [
        ("登録した瞬間から業界知識80%プリロード済み",
         "建設・製造・医療福祉など業種別テンプレートが最初から入っている。\nゼロから学習させる必要なし。初日から価値が出る。"),
        ("自社データが増えるほど精度が上がる",
         "見積書・議事録・報告書をアップロードするだけ。\nAIが自社固有の知識・判断基準を学習し続ける。"),
        ("業界全体が賢くなる（匿名ベンチマーク）",
         "他社の匿名集合知が全ユーザーにフィードバック。\n「同業他社の積算単価」「業界平均の不良率」が参照可能。"),
    ]

    top_positions = [Cm(5.2), Cm(8.8), Cm(12.4)]

    for i, (title, desc) in enumerate(points):
        top = top_positions[i]
        add_rect(slide, Cm(1.5), top, Cm(0.4), Cm(2.8),
                 RED if i == 0 else RGBColor(0x44, 0x66, 0x99))
        add_textbox(
            slide, Cm(2.5), top + Cm(0.1),
            SLIDE_WIDTH - Cm(4), Cm(1.0),
            title, font_size=16, bold=True, color=WHITE
        )
        add_textbox(
            slide, Cm(2.5), top + Cm(1.1),
            SLIDE_WIDTH - Cm(4), Cm(1.5),
            desc, font_size=12, color=RGBColor(0xBB, 0xCC, 0xEE)
        )

    add_slide_number(slide, 4, text_color=RGBColor(0x55, 0x66, 0x88))


def slide_05_input(prs):
    """Slide 5: メール/Excel で渡すだけ（LINE は Coming Soon）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, RGBColor(0xF8, 0xF9, 0xFA))

    add_section_title(slide, "メール / Excel で渡すだけ",
                      subtitle="操作不要。社長の使い慣れたツールからそのまま入力できる")
    add_divider(slide, Cm(2.3))

    # ---- カード定義 ----
    # (title, bullets, accent_color, is_coming_soon)
    channels = [
        ("メール",
         ["メールを転送するだけ", "議事録・報告書も自動解析", "返信下書きも自動生成"],
         ACCENT_BLUE,
         False),
        ("Excel / PDF",
         ["ファイルをアップロードするだけ", "見積書・請求書・図面を自動読み取り", "OCR対応（手書き帳票もOK）"],
         RGBColor(0x0F, 0x7B, 0x6C),
         False),
        ("LINE",
         ["LINE で送信するだけ", "音声メッセージもOK", "承認・却下も LINE で完結"],
         RGBColor(0xA0, 0xA8, 0xB4),
         True),
    ]

    card_w = Cm(9.8)
    card_h = Cm(11.8)
    gap = Cm(1.4)
    total_w = card_w * 3 + gap * 2
    start_left = (SLIDE_WIDTH - total_w) / 2
    top = Cm(3.0)

    CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
    CARD_BG_DIM = RGBColor(0xEA, 0xEC, 0xEF)
    TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
    TEXT_DIM    = RGBColor(0x9A, 0xA0, 0xAC)
    BULLET_OK   = RGBColor(0x34, 0xA8, 0x53)

    for i, (title, bullets, accent, coming_soon) in enumerate(channels):
        left = start_left + i * (card_w + gap)
        bg = CARD_BG_DIM if coming_soon else CARD_BG

        # カード本体（白い角丸風 — pptx は角丸矩形 shape type 5）
        card = slide.shapes.add_shape(5, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = bg
        card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
        card.line.width = Pt(0.75)

        # アクセントバー（上部）
        bar_color = RGBColor(0xC8, 0xCE, 0xD8) if coming_soon else accent
        bar = slide.shapes.add_shape(1, left, top, card_w, Cm(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_color
        bar.line.fill.background()

        # アイコン代わりの丸（アクセントカラー）
        icon_size = Cm(1.4)
        icon_left = left + (card_w - icon_size) / 2
        ic = slide.shapes.add_shape(9, icon_left, top + Cm(0.8), icon_size, icon_size)  # 9=oval
        ic.fill.solid()
        ic.fill.fore_color.rgb = RGBColor(0xC8, 0xCE, 0xD8) if coming_soon else accent
        ic.line.fill.background()

        # タイトル
        title_color = TEXT_DIM if coming_soon else TEXT_DARK
        tb_title = slide.shapes.add_textbox(
            left + Cm(0.6), top + Cm(2.6), card_w - Cm(1.2), Cm(1.3))
        tf = tb_title.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.name = FONT_NAME
        run.font.color.rgb = title_color

        # 区切り線
        sep = slide.shapes.add_connector(
            1,
            left + Cm(1.2), top + Cm(4.1),
            left + card_w - Cm(1.2), top + Cm(4.1)
        )
        sep.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
        sep.line.width = Pt(0.5)

        # 箇条書き（bullet point）
        for j, bullet in enumerate(bullets):
            b_top = top + Cm(4.5) + Cm(j * 2.0)
            b_color = TEXT_DIM if coming_soon else TEXT_DARK

            # チェックマーク
            ck = slide.shapes.add_textbox(
                left + Cm(0.7), b_top, Cm(0.7), Cm(0.9))
            ck_tf = ck.text_frame
            ck_p = ck_tf.paragraphs[0]
            ck_run = ck_p.add_run()
            ck_run.text = "—" if coming_soon else "✓"
            ck_run.font.size = Pt(11)
            ck_run.font.bold = True
            ck_run.font.name = "Arial"
            ck_run.font.color.rgb = TEXT_DIM if coming_soon else BULLET_OK

            # テキスト
            bt = slide.shapes.add_textbox(
                left + Cm(1.5), b_top, card_w - Cm(2.0), Cm(0.9))
            bt_tf = bt.text_frame
            bt_p = bt_tf.paragraphs[0]
            bt_run = bt_p.add_run()
            bt_run.text = bullet
            bt_run.font.size = Pt(12)
            bt_run.font.name = FONT_NAME
            bt_run.font.color.rgb = b_color

        # COMING SOON バッジ（LINE のみ）
        if coming_soon:
            badge_w = Cm(5.5)
            badge_h = Cm(0.9)
            badge_left = left + (card_w - badge_w) / 2
            badge_top = top + card_h - Cm(1.8)
            badge = slide.shapes.add_shape(5, badge_left, badge_top, badge_w, badge_h)
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(0x5F, 0x67, 0x7A)
            badge.line.fill.background()

            tb_cs = slide.shapes.add_textbox(
                badge_left, badge_top + Cm(0.1), badge_w, badge_h)
            cs_tf = tb_cs.text_frame
            cs_p = cs_tf.paragraphs[0]
            cs_p.alignment = PP_ALIGN.CENTER
            cs_run = cs_p.add_run()
            cs_run.text = "COMING SOON"
            cs_run.font.size = Pt(10)
            cs_run.font.bold = True
            cs_run.font.name = "Arial"
            cs_run.font.color.rgb = WHITE

    # フッターメッセージ
    add_textbox(
        slide,
        Cm(2.0), Cm(15.4),
        SLIDE_WIDTH - Cm(4), Cm(1.2),
        "新しいツールの導入・社員教育は不要。今日から使い始められる。",
        font_size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, 5)


def slide_06_day1(prs):
    """Slide 6: Day 1 Value"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "Day 1 Value",
                      subtitle="登録した初日から、3つの機能が使える")
    add_divider(slide, Cm(2.3))

    features = [
        ("Q&A（24時間対応）",
         "社長に聞かなくて良い\n\n「杭工事の坪単価は？」\n「この契約条件は問題ない？」\n→ 業界テンプレートで即回答",
         "業界テンプレート\n80%プリロード済み"),
        ("テンプレート自動生成",
         "書き方に迷わない\n\n見積書・契約書・報告書・議事録\n→ ひな形を一発生成\n→ 自社フォーマットに自動適応",
         "50種類以上の\nテンプレート標準搭載"),
        ("書類作成・自動入力",
         "転記・コピペ不要\n\n請求書 → 帳簿に自動転記\nメール → 案件管理に自動登録\n名刺 → 顧客DBに自動入力",
         "API連携で\nkintone/freeeも自動化"),
    ]

    positions = [Cm(1.5), Cm(12.3), Cm(23.1)]
    box_width = Cm(10.0)

    for i, (title, body, badge) in enumerate(features):
        left = positions[i]
        top = Cm(3.2)

        add_rect(slide, left, top, box_width, Cm(10.5), LIGHT_GRAY)
        add_rect(slide, left, top, box_width, Cm(1.8), NAVY)

        add_textbox(
            slide, left + Cm(0.4), top + Cm(0.4),
            box_width - Cm(0.8), Cm(1.2),
            title, font_size=16, bold=True, color=WHITE
        )
        add_textbox(
            slide, left + Cm(0.4), top + Cm(2.2),
            box_width - Cm(0.8), Cm(6.0),
            body, font_size=12, color=DARK_NAVY
        )
        # バッジ
        add_rect(slide, left + Cm(0.4), top + Cm(8.5),
                 box_width - Cm(0.8), Cm(1.5),
                 RED)
        add_textbox(
            slide, left + Cm(0.4), top + Cm(8.6),
            box_width - Cm(0.8), Cm(1.3),
            badge, font_size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

    add_slide_number(slide, 6)


def slide_07_demo(prs):
    """Slide 7: デモ画面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, DARK_NAVY)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Cm(0.4), RED)

    add_textbox(
        slide, Cm(2), Cm(1.2),
        SLIDE_WIDTH - Cm(4), Cm(1.5),
        "ここでライブデモ",
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    # デモ画面枠
    add_rect(slide, Cm(2), Cm(3.0), SLIDE_WIDTH - Cm(4), Cm(11.5),
             RGBColor(0x1A, 0x2A, 0x4A))
    # 枠線
    demo_box = slide.shapes.add_shape(
        1, Cm(2), Cm(3.0), SLIDE_WIDTH - Cm(4), Cm(11.5)
    )
    demo_box.fill.background()
    demo_box.line.color.rgb = RGBColor(0x44, 0x66, 0x99)
    demo_box.line.width = Pt(2)

    add_textbox(
        slide, Cm(2), Cm(7.5),
        SLIDE_WIDTH - Cm(4), Cm(2.5),
        "【 画面キャプチャ / ライブデモ 】\n\nQ&A・BPO自動化・能動提案の3機能をデモ",
        font_size=16, color=RGBColor(0x77, 0x99, 0xCC),
        align=PP_ALIGN.CENTER
    )

    # デモポイント
    demo_points = ["① 社員がLINEで質問 → 即回答",
                   "② 見積書アップロード → 自動解析",
                   "③ 「来月の人員不足を検知」通知"]
    add_textbox(
        slide, Cm(2.5), Cm(15.0),
        SLIDE_WIDTH - Cm(5), Cm(2.5),
        "  ".join(demo_points),
        font_size=12, color=RGBColor(0xAA, 0xBB, 0xDD)
    )

    add_slide_number(slide, 7, text_color=RGBColor(0x55, 0x66, 0x88))


def slide_08_cases(prs):
    """Slide 8: 業種別事例"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "業種別 導入事例",
                      subtitle="業界テンプレートで初日から効果が出る")
    add_divider(slide, Cm(2.3))

    cases = [
        ("建設業",
         "三栄建設（210名・売上98億）\n\n積算時間\n▼ 60%削減\n120h → 48h/月\n\n・先代の積算ノウハウをAIが継承\n・協力会社評価の属人化を解消\n・受注判断をAIが補助",
         NAVY, "ROI 2,701%"),
        ("製造業",
         "大和精密工業（145名・売上52億）\n\n不良率\n▼ 46%改善\n2.8% → 1.5%\n\n・加工条件の暗黙知を形式知化\n・新人育成期間 12ヶ月→6ヶ月\n・品質判断ナレッジの永続化",
         LIGHT_NAVY, "ROI 2,239%"),
        ("医療・福祉",
         "メディカルスマイルG（20名・3院）\n\n記録業務\n▼ 50%削減\n\n・カウンセリング成約率 +15%\n・スタッフ教育コスト -40%\n・カウンセリングスクリプト自動生成",
         RGBColor(0x0A, 0x52, 0x94), "ROI 5,227%"),
    ]

    positions = [Cm(1.5), Cm(12.3), Cm(23.1)]
    box_width = Cm(10.0)

    for i, (industry, body, color, roi) in enumerate(cases):
        left = positions[i]
        top = Cm(3.2)
        box_h = Cm(10.8)

        add_rect(slide, left, top, box_width, box_h, color)
        add_textbox(
            slide, left + Cm(0.4), top + Cm(0.3),
            box_width - Cm(0.8), Cm(1.0),
            industry, font_size=18, bold=True, color=WHITE
        )
        add_textbox(
            slide, left + Cm(0.4), top + Cm(1.5),
            box_width - Cm(0.8), Cm(7.5),
            body, font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF)
        )
        # ROIバッジ
        add_rect(slide, left + Cm(0.4), top + box_h - Cm(1.8),
                 box_width - Cm(0.8), Cm(1.4), RED)
        add_textbox(
            slide, left + Cm(0.4), top + box_h - Cm(1.75),
            box_width - Cm(0.8), Cm(1.2),
            roi, font_size=14, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

    add_slide_number(slide, 8)


def slide_09_twin(prs):
    """Slide 9: 9次元モデル・能動提案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "9次元モデル × 能動提案",
                      subtitle="AIが先にリスクを検知して社長に通知する")
    add_divider(slide, Cm(2.3))

    # 左側: 5次元モデル説明
    add_textbox(
        slide, Cm(1.5), Cm(3.2),
        Cm(14), Cm(1.2),
        "デジタルツイン 5次元（Phase 1）",
        font_size=15, bold=True, color=NAVY
    )

    dimensions = [
        ("① ヒト", "スキル・評価・退職リスク"),
        ("② プロセス", "業務フロー・ボトルネック"),
        ("③ コスト", "原価・人件費・キャッシュ"),
        ("④ ツール", "使用SaaS・稼働率"),
        ("⑤ リスク", "コンプライアンス・品質"),
    ]

    for j, (dim, desc) in enumerate(dimensions):
        top = Cm(4.5) + j * Cm(2.2)
        add_rect(slide, Cm(1.5), top, Cm(3.2), Cm(1.8),
                 NAVY if j == 0 else LIGHT_GRAY)
        add_textbox(
            slide, Cm(1.6), top + Cm(0.3),
            Cm(3.0), Cm(1.2),
            dim, font_size=13, bold=True,
            color=WHITE if j == 0 else NAVY,
            align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, Cm(5.2), top + Cm(0.3),
            Cm(10.0), Cm(1.2),
            desc, font_size=12, color=DARK_NAVY
        )

    # 右側: 能動提案の例
    add_rect(slide, Cm(17), Cm(3.2), Cm(14.5), Cm(12.5),
             RGBColor(0xF0, 0xF4, 0xFF))
    add_textbox(
        slide, Cm(17.5), Cm(3.5),
        Cm(13.5), Cm(1.2),
        "能動提案の例",
        font_size=15, bold=True, color=NAVY
    )

    proactive_examples = [
        "📊 「来月、A現場で人員が3名不足する見込みです」",
        "⚠️  「B協力会社の評価スコアが低下しています」",
        "💰 「今期の粗利率が業界平均を下回っています」",
        "📋 「C社の契約更新期限が30日後に迫っています」",
        "🔍 「同業他社の見積単価より15%高い傾向にあります」",
    ]

    for k, example in enumerate(proactive_examples):
        top = Cm(5.0) + k * Cm(1.9)
        add_rect(slide, Cm(17.5), top, Cm(13.0), Cm(1.6),
                 WHITE)
        add_textbox(
            slide, Cm(18.0), top + Cm(0.2),
            Cm(12.5), Cm(1.2),
            example, font_size=11, color=DARK_NAVY
        )

    add_slide_number(slide, 9)


def slide_10_bpo(prs):
    """Slide 10: BPO自動化レベル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "BPO自動化",
                      subtitle="使うほど自動化レベルが上がる — 5段階の進化")
    add_divider(slide, Cm(2.3))

    levels = [
        ("Lv.1", "情報収集",
         "Q&A・書類検索・テンプレート提示", LIGHT_GRAY, NAVY),
        ("Lv.2", "提案・下書き",
         "見積下書き・報告書自動生成", LIGHT_GRAY, NAVY),
        ("Lv.3", "確認承認",
         "LINE承認・例外検知・差し戻し", RGBColor(0xE8, 0xF0, 0xFF), NAVY),
        ("Lv.4", "自動実行",
         "請求書発行・入金消込・API連携", NAVY, WHITE),
        ("Lv.5", "完全自律",
         "予測発注・人員最適化・自律改善", RED, WHITE),
    ]

    box_width = Cm(5.8)
    for i, (lv, title, desc, bg, text_color) in enumerate(levels):
        left = Cm(1.5) + i * Cm(6.3)
        top = Cm(3.3)
        box_h = Cm(10.0)

        add_rect(slide, left, top, box_width, box_h, bg)

        add_textbox(
            slide, left + Cm(0.3), top + Cm(0.5),
            box_width - Cm(0.6), Cm(1.0),
            lv, font_size=20, bold=True, color=RED if bg != RED else WHITE,
            align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, left + Cm(0.3), top + Cm(1.7),
            box_width - Cm(0.6), Cm(1.2),
            title, font_size=15, bold=True, color=text_color,
            align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, left + Cm(0.3), top + Cm(3.2),
            box_width - Cm(0.6), Cm(5.5),
            desc, font_size=11, color=text_color
        )

    # 矢印テキスト
    add_textbox(
        slide, Cm(1.5), Cm(14.2),
        SLIDE_WIDTH - Cm(3), Cm(1.0),
        "← 今すぐ使える  ─────────────────────────────  使い続けると到達 →",
        font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, 10)


def slide_11_pricing(prs):
    """Slide 11: 料金プラン"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "料金プラン",
                      subtitle="モジュール選択制 — 必要なものだけON/OFFできる")
    add_divider(slide, Cm(2.3))

    plans = [
        ("ブレイン",
         "¥30,000",
         "/月",
         "ナレッジ管理\nQ&A（500回/月）\nデジタルツイン\nテンプレート生成\n書類作成補助",
         "全業界共通\nまずここから",
         LIGHT_GRAY, NAVY, False),
        ("BPOコア",
         "¥250,000",
         "/月",
         "業種特化BPOモジュール群\n主要業務の自動化\nBPO処理（300回/月）\nAPI連携（kintone/freee）\nLINE承認フロー",
         "建設/製造/医療福祉\n不動産/物流/卸売",
         NAVY, WHITE, True),
        ("建設BPO",
         "¥300,000",
         "/月",
         "積算自動化\n協力会社管理\n工程管理AIアシスト\n書類電子化\n現場報告AI",
         "業種特化プレミアム\n建設業フルパック",
         LIGHT_GRAY, NAVY, False),
    ]

    positions = [Cm(1.5), Cm(12.3), Cm(23.1)]
    box_width = Cm(10.0)

    for i, (name, price, unit, features, badge, bg, text_color, is_featured) in enumerate(plans):
        left = positions[i]
        top = Cm(3.1)
        box_h = Cm(12.5)

        add_rect(slide, left, top, box_width, box_h, bg)
        if is_featured:
            add_rect(slide, left, top, box_width, Cm(0.35), RED)

        # プラン名
        add_textbox(
            slide, left + Cm(0.5), top + Cm(0.5),
            box_width - Cm(1), Cm(1.0),
            name, font_size=18, bold=True, color=text_color
        )
        # 価格
        add_textbox(
            slide, left + Cm(0.5), top + Cm(1.7),
            box_width - Cm(1), Cm(1.5),
            price, font_size=28, bold=True, color=RED,
            align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, left + Cm(0.5), top + Cm(3.0),
            box_width - Cm(1), Cm(0.8),
            unit, font_size=12,
            color=text_color, align=PP_ALIGN.CENTER
        )

        # 機能リスト
        add_textbox(
            slide, left + Cm(0.5), top + Cm(4.0),
            box_width - Cm(1), Cm(6.5),
            features, font_size=11, color=text_color
        )

        # バッジ
        badge_bg = RED if not is_featured else RGBColor(0x44, 0x66, 0x99)
        add_rect(slide, left + Cm(0.5), top + box_h - Cm(1.8),
                 box_width - Cm(1), Cm(1.4), badge_bg)
        add_textbox(
            slide, left + Cm(0.5), top + box_h - Cm(1.7),
            box_width - Cm(1), Cm(1.2),
            badge, font_size=10, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

    add_slide_number(slide, 11)


def slide_12_roi(prs):
    """Slide 12: ROI計算"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Cm(0.4), RED)

    add_textbox(
        slide, Cm(2), Cm(1.0),
        SLIDE_WIDTH - Cm(4), Cm(1.8),
        "ROI シミュレーション（建設業・30名）",
        font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    # 投資 vs 効果
    # 左: 投資
    add_rect(slide, Cm(2), Cm(3.2), Cm(13.5), Cm(9.0),
             RGBColor(0x1A, 0x2A, 0x4A))
    add_textbox(
        slide, Cm(2.5), Cm(3.5),
        Cm(12.5), Cm(1.2),
        "投資（月額）",
        font_size=16, bold=True, color=RGBColor(0xAA, 0xBB, 0xDD)
    )
    add_textbox(
        slide, Cm(2.5), Cm(5.0),
        Cm(12.5), Cm(2.5),
        "¥360,000 / 月",
        font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    invest_items = "ブレイン ¥30,000\n建設BPO ¥300,000\n（ブレイン+BPO合計）"
    add_textbox(
        slide, Cm(2.5), Cm(7.8),
        Cm(12.5), Cm(3.0),
        invest_items, font_size=12,
        color=RGBColor(0xBB, 0xCC, 0xEE)
    )

    # 中央: 矢印
    add_textbox(
        slide, Cm(15.5), Cm(6.5),
        Cm(3.0), Cm(2.0),
        "→",
        font_size=36, bold=True, color=RED, align=PP_ALIGN.CENTER
    )

    # 右: 効果
    add_rect(slide, Cm(18.5), Cm(3.2), Cm(13.5), Cm(9.0),
             RGBColor(0x1A, 0x2A, 0x4A))
    add_textbox(
        slide, Cm(19.0), Cm(3.5),
        Cm(12.5), Cm(1.2),
        "効果（月額換算）",
        font_size=16, bold=True, color=RGBColor(0xAA, 0xBB, 0xDD)
    )
    add_textbox(
        slide, Cm(19.0), Cm(5.0),
        Cm(12.5), Cm(2.5),
        "¥10,080,000 / 月",
        font_size=24, bold=True, color=RED, align=PP_ALIGN.CENTER
    )
    effect_items = "積算時間削減: ¥3,888,000\n属人化リスク解消: ¥1,944,000\n業務自動化: ¥3,024,000\nその他効率化: ¥1,224,000"
    add_textbox(
        slide, Cm(19.0), Cm(7.8),
        Cm(12.5), Cm(3.0),
        effect_items, font_size=11,
        color=RGBColor(0xBB, 0xCC, 0xEE)
    )

    # ROI
    add_rect(slide, Cm(2), Cm(13.0), SLIDE_WIDTH - Cm(4), Cm(3.5), RED)
    add_textbox(
        slide, Cm(2), Cm(13.2),
        SLIDE_WIDTH - Cm(4), Cm(3.0),
        "ROI  2,701%",
        font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, 12, text_color=RGBColor(0x55, 0x66, 0x88))


def slide_13_security(prs):
    """Slide 13: セキュリティ&データ所在"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "セキュリティ & データ所在",
                      subtitle="freee・SmartHRと同じ構造。国内クラウド・銀行水準の暗号化")
    add_divider(slide, Cm(2.3))

    sec_items = [
        ("日本国内保管",
         "DB・アプリ・ファイル全て国内\n東京データセンター（AWS/GCP）\nfreee・kintoneと同じ構造"),
        ("暗号化",
         "保存: AES-256-GCM（銀行水準）\n通信: TLS 1.3（最新規格）\n個人情報は自動マスキング後に送信"),
        ("テナント完全分離",
         "他社データとは壁で完全隔離\nRow Level Security（RLS）\n御社のデータが他社に漏れない"),
        ("AI学習に使用しない",
         "御社データをAI学習に使わない\nGoogle DPA（処理委託契約）締結済\n処理後はGoogleに保存されない"),
        ("監査ログ",
         "誰が・何を・いつしたか全記録\n改ざん不可な監査証跡\nエンタープライズ向けSSOも対応"),
        ("解約時完全削除",
         "全データを完全削除\n削除証明書を発行\nいつでも解約可（違約金なし）"),
    ]

    positions = [(Cm(1.5), Cm(3.2)), (Cm(12.3), Cm(3.2)), (Cm(23.1), Cm(3.2)),
                 (Cm(1.5), Cm(10.5)), (Cm(12.3), Cm(10.5)), (Cm(23.1), Cm(10.5))]
    box_width = Cm(10.0)
    box_h = Cm(6.5)

    for i, (title, desc) in enumerate(sec_items):
        left, top = positions[i]
        add_rect(slide, left, top, box_width, box_h, LIGHT_GRAY)
        add_rect(slide, left, top, Cm(0.3), box_h, NAVY)
        add_textbox(
            slide, left + Cm(0.7), top + Cm(0.3),
            box_width - Cm(1), Cm(1.0),
            title, font_size=15, bold=True, color=NAVY
        )
        add_textbox(
            slide, left + Cm(0.7), top + Cm(1.5),
            box_width - Cm(1), Cm(4.5),
            desc, font_size=11, color=DARK_NAVY
        )

    add_slide_number(slide, 13)


def slide_14_flow(prs):
    """Slide 14: 導入フロー"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)

    add_section_title(slide, "導入フロー",
                      subtitle="招待制登録 → ナレッジ入力 → Q&A開始 → BPO設定。最短3ステップ")
    add_divider(slide, Cm(2.3))

    steps = [
        ("STEP 1", "招待制登録", "招待URLから\n30秒で登録完了\n（パスワード設定のみ）",
         "Day 0"),
        ("STEP 2", "ナレッジ入力", "LINE/メール/ファイルで\n会社情報を入力\n（5分〜）",
         "Day 1"),
        ("STEP 3", "Q&A開始", "業界テンプレートで\n即日Q&A利用開始\n（Aha Moment）",
         "Day 1"),
        ("STEP 4", "BPO設定", "自動化したい業務を\n1つ選んでON\n（1-2時間）",
         "Week 1"),
    ]

    box_width = Cm(7.0)
    for i, (step, title, desc, timing) in enumerate(steps):
        left = Cm(1.5) + i * Cm(7.8)
        top = Cm(3.2)
        box_h = Cm(10.5)

        add_rect(slide, left, top, box_width, box_h, LIGHT_GRAY)
        add_rect(slide, left, top, box_width, Cm(2.5), NAVY)

        add_textbox(
            slide, left + Cm(0.3), top + Cm(0.2),
            box_width - Cm(0.6), Cm(0.9),
            step, font_size=13, bold=True, color=RED
        )
        add_textbox(
            slide, left + Cm(0.3), top + Cm(1.1),
            box_width - Cm(0.6), Cm(1.1),
            title, font_size=16, bold=True, color=WHITE
        )
        add_textbox(
            slide, left + Cm(0.3), top + Cm(3.0),
            box_width - Cm(0.6), Cm(5.5),
            desc, font_size=13, color=DARK_NAVY
        )
        # タイミングバッジ
        add_rect(slide, left + Cm(0.3), top + box_h - Cm(1.6),
                 box_width - Cm(0.6), Cm(1.3), RED)
        add_textbox(
            slide, left + Cm(0.3), top + box_h - Cm(1.5),
            box_width - Cm(0.6), Cm(1.1),
            timing, font_size=13, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

        # 矢印
        if i < 3:
            add_textbox(
                slide, left + box_width + Cm(0.1), top + Cm(3.5),
                Cm(0.7), Cm(2.0),
                "▶", font_size=14, bold=True, color=NAVY
            )

    add_slide_number(slide, 14)


def slide_15_cta(prs):
    """Slide 15: CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, 0, 0, SLIDE_WIDTH, Cm(0.5), RED)
    add_rect(slide, 0, SLIDE_HEIGHT - Cm(0.5), SLIDE_WIDTH, Cm(0.5), RED)

    # メインCTA
    add_textbox(
        slide, Cm(2), Cm(2.5),
        SLIDE_WIDTH - Cm(4), Cm(2.0),
        "まずブレインで始めてみる",
        font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Cm(2), Cm(4.8),
        SLIDE_WIDTH - Cm(4), Cm(1.2),
        "月¥30,000 から。今すぐ招待URLを発行します。",
        font_size=18, color=RGBColor(0xCC, 0xDD, 0xFF),
        align=PP_ALIGN.CENTER
    )

    # アクションボックス
    add_rect(slide, Cm(5), Cm(6.5), SLIDE_WIDTH - Cm(10), Cm(5.5), RED)
    add_textbox(
        slide, Cm(5), Cm(7.0),
        SLIDE_WIDTH - Cm(10), Cm(4.5),
        "▶  招待URLを発行する\n\n今日中に invite@shachotwo.jp へご連絡ください\n または営業担当者に直接お申し付けください",
        font_size=14, color=WHITE, align=PP_ALIGN.CENTER
    )

    # 3つのポイント
    points = [
        "初期費用 ゼロ",
        "いつでも解約可",
        "7日間 無料トライアル"
    ]
    pos = [Cm(2), Cm(12.5), Cm(23.5)]
    for k, (point, left) in enumerate(zip(points, pos)):
        add_rect(slide, left, Cm(13.5), Cm(9.5), Cm(2.5),
                 RGBColor(0x1A, 0x2A, 0x4A))
        add_textbox(
            slide, left, Cm(13.7),
            Cm(9.5), Cm(2.0),
            point, font_size=14, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER
        )

    add_textbox(
        slide, Cm(2), Cm(17.0),
        SLIDE_WIDTH - Cm(4), Cm(0.9),
        "© 2026 シャチョツー（湊合同会社）  |  invite@shachotwo.jp",
        font_size=10, color=MID_GRAY, align=PP_ALIGN.CENTER
    )

    add_slide_number(slide, 15, text_color=RGBColor(0x55, 0x66, 0x88))


# =========================================
# メイン処理
# =========================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating slides...")

    slide_01_cover(prs)
    print("  Slide 1: 表紙")

    slide_02_problem(prs)
    print("  Slide 2: こんなことありませんか？")

    slide_03_what(prs)
    print("  Slide 3: シャチョツーとは")

    slide_04_network(prs)
    print("  Slide 4: データネットワーク効果")

    slide_05_input(prs)
    print("  Slide 5: LINE/メール/Excelで渡すだけ")

    slide_06_day1(prs)
    print("  Slide 6: Day 1 Value")

    slide_07_demo(prs)
    print("  Slide 7: デモ画面")

    slide_08_cases(prs)
    print("  Slide 8: 業種別事例")

    slide_09_twin(prs)
    print("  Slide 9: 9次元モデル・能動提案")

    slide_10_bpo(prs)
    print("  Slide 10: BPO自動化")

    slide_11_pricing(prs)
    print("  Slide 11: 料金プラン")

    slide_12_roi(prs)
    print("  Slide 12: ROI計算")

    slide_13_security(prs)
    print("  Slide 13: セキュリティ&データ所在")

    slide_14_flow(prs)
    print("  Slide 14: 導入フロー")

    slide_15_cta(prs)
    print("  Slide 15: CTA")

    output_path = "/Users/sugimotoyuuhi/code/ai_agent/shachotwo/c_事業計画/シャチョツー_営業デモ資料.pptx"
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print("Done!")


if __name__ == "__main__":
    main()
