"""
シャチョツー 営業スライド生成スクリプト（3業種統合版）

使い方:
  python shachotwo/generate_slides.py

出力:
  shachotwo/shachotwo_deck.pptx（建設業・製造業・歯科 統合）

Canvaへの取り込み:
  1. https://www.canva.com を開く
  2. ホーム右上の「アップロード」をクリック
  3. pptx をドラッグ＆ドロップ → 自動でCanvaプレゼンに変換
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 色定数 ──
NAVY = RGBColor(0x1B, 0x36, 0x5D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)

# 業種カラー
IND_COLORS = {
    "建設業": RGBColor(0xD4, 0x6B, 0x08),  # オレンジ系
    "製造業": RGBColor(0x1B, 0x7A, 0x3D),  # グリーン系
    "歯科":   RGBColor(0x6C, 0x3A, 0xB5),  # パープル系
}


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ヘルパー関数
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=DARK, bold=False,
                  space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    return p


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_table(slide, data, left, top, width, height):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    table = shape.table
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14 if cols >= 4 else 15)
                p.font.color.rgb = WHITE if i == 0 else DARK
                p.font.bold = (i == 0 or j == 0)
                if j >= 1 and i > 0 and val.startswith("○"):
                    p.font.color.rgb = ACCENT
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
    return table


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド生成（統合版）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def generate_unified_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 表紙 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_textbox(s, 1, 1.2, 11, 1.5,
                "シャチョツー（社長2号）",
                font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1, 3.0, 11, 0.8,
                "会社のNo.2を身近に。手軽に。",
                font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1, 4.5, 11, 0.5,
                "建設業 ・ 製造業 ・ 歯科  対応",
                font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: 課題（3業種共通） ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "こんな課題ありませんか？",
                font_size=36, color=NAVY, bold=True)
    problems = [
        ("社長/院長にしかわからないことが多すぎる",
         "出張中・不在時に判断が止まる → 現場・スタッフが動けない"),
        ("ベテランが辞めたらノウハウが消える",
         "暗黙知の引継ぎは「見て覚えろ」。新人の独り立ちに半年〜1年以上"),
        ("情報がバラバラで見つからない",
         "紙の帳票、Excel、レセコン、LINEのやり取り…\n「あの情報どこだっけ」で毎日30分消える"),
        ("同じ質問に何度も答えている",
         "社員/スタッフからの問い合わせが社長/院長に集中。月20-30時間を消費"),
        ("現場の問題に気づくのがいつも遅い",
         "数字を見て「やばい」と思った時にはもう手遅れ"),
    ]
    for i, (title, desc) in enumerate(problems):
        y = 1.5 + i * 1.1
        tf = add_textbox(s, 1.2, y, 11, 0.5, f"✕  {title}",
                         font_size=22, color=DARK, bold=True)
        add_paragraph(tf, f"    {desc}", font_size=16, color=GRAY)

    # ── Slide 3: シャチョツーとは ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "シャチョツーとは", font_size=36, color=NAVY, bold=True)
    add_textbox(s, 1.5, 1.8, 10, 1,
                "シャチョツー ＝ 社長の「黒子AI」",
                font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1.5, 3.0, 10, 1.5,
                "社長が話すだけで、会社のナレッジが整理され、\n"
                "社員が質問でき、AIがリスクを先に教えてくれる。",
                font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 1.5, 5.0, 10, 1,
                "社長が話す  →  AIが構造化  →  社員が活用  →  AIが提案",
                font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 4: 3つの価値 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "3つの価値", font_size=36, color=NAVY, bold=True)
    values = [
        ("❶ 話すだけでナレッジが貯まる",
         "テキスト・音声・書類、なんでもOK。\nAIが自動で整理・分類・構造化。",
         "建設: 積算基準、協力会社評価\n製造: 加工条件、品質基準\n歯科: 診療方針、自費案内基準"),
        ("❷ 社員がAIに聞ける",
         "引用付きで即回答。社長の時間を奪わない。",
         "建設: 「この案件の積算基準は？」\n製造: 「この製品の品質基準は？」\n歯科: 「この処置の院内基準は？」"),
        ("❸ AIが先にリスクを教える",
         "社長が気づく前にAIが気づく。",
         "建設: 「○○工事の工期が危ない」\n製造: 「不良率が先月比1.5倍です」\n歯科: 「リコール率が目標を下回っています」"),
    ]
    for i, (title, desc, examples) in enumerate(values):
        x = 0.8 + i * 4.1
        tf = add_textbox(s, x, 1.6, 3.8, 0.6, title,
                         font_size=22, color=DARK, bold=True)
        add_paragraph(tf, desc, font_size=17, color=GRAY)
        add_paragraph(tf, "", font_size=8)
        add_paragraph(tf, examples, font_size=14, color=ACCENT)

    # ── Slide 5: 業界テンプレート一覧 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "業界テンプレートで80%完成スタート",
                font_size=36, color=NAVY, bold=True)
    add_textbox(s, 0.8, 1.3, 11, 0.6,
                "「うちの業界わかるの？」→ わかります。導入初日から80%完成。",
                font_size=22, color=DARK, bold=True)

    templates = [
        ("建設業", IND_COLORS["建設業"], [
            "部署: 営業 / 工事 / 品質管理 / 経理 / 総務",
            "フロー: 受注→施工計画→施工→検査→引渡→請求",
            "ルール: 積算基準 / 原価管理 / 安全管理",
            "課題: 紙帳票 / Excel属人化 / 二重入力",
        ]),
        ("製造業", IND_COLORS["製造業"], [
            "部署: 営業 / 製造 / 品質管理 / 購買 / 経理",
            "フロー: 受注→生産計画→製造→検品→出荷→請求",
            "ルール: 原価計算 / 在庫管理 / 品質基準(ISO)",
            "課題: 紙帳票 / 在庫の見える化 / 品質記録",
        ]),
        ("歯科", IND_COLORS["歯科"], [
            "部門: 受付 / 診療 / 歯科衛生士 / 技工 / 経理",
            "フロー: 予約→受付→診療→会計→リコール管理",
            "ルール: 自費/保険案内基準 / リコール間隔",
            "課題: レセコン二重管理 / リコール漏れ / 自費率",
        ]),
    ]
    for i, (name, ind_color, items) in enumerate(templates):
        x = 0.6 + i * 4.2
        tf = add_textbox(s, x, 2.2, 3.9, 0.5, f"■ {name}",
                         font_size=22, color=ind_color, bold=True)
        for item in items:
            add_paragraph(tf, f"  {item}", font_size=15, color=GRAY, space_before=Pt(8))

    # ── Slide 6: 建設業の活用シーン ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "活用シーン ― 建設業",
                font_size=36, color=IND_COLORS["建設業"], bold=True)
    add_textbox(s, 0.8, 1.3, 6, 0.5,
                "社長の1日（導入後）",
                font_size=24, color=DARK, bold=True)
    timeline_c = [
        ("朝（5分）", "ダッシュボードでAI提案を確認\n「△△工事の協力会社手配が\n遅れています」"),
        ("日中", "社員の質問はAIが回答\n社長への問い合わせ激減"),
        ("移動中", "音声入力で話すだけ\nナレッジに自動追加"),
        ("夕方（3分）", "AIの提案をレビュー\n承認 or 却下をタップ"),
    ]
    for i, (time, desc) in enumerate(timeline_c):
        x = 0.6 + i * 3.1
        add_textbox(s, x, 2.2, 2.8, 0.5, time,
                    font_size=20, color=IND_COLORS["建設業"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.9, 2.8, 2.0, desc,
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(s, 0.8, 5.3, 11, 0.5,
                     "Q&A例: 「この案件の積算単価はどのくらいが妥当？」",
                     font_size=18, color=DARK, bold=True)
    add_paragraph(tf,
                  "→ 「過去3年の類似案件は坪単価¥85,000。地盤改良込みで¥92,000-95,000が妥当です」",
                  font_size=17, color=ACCENT)

    # ── Slide 7: 製造業の活用シーン ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "活用シーン ― 製造業",
                font_size=36, color=IND_COLORS["製造業"], bold=True)
    add_textbox(s, 0.8, 1.3, 6, 0.5,
                "社長の1日（導入後）",
                font_size=24, color=DARK, bold=True)
    timeline_m = [
        ("朝（5分）", "ダッシュボードでAI提案を確認\n「A製品の不良率が\n先月比1.5倍です」"),
        ("日中", "社員の質問はAIが回答\n社長への問い合わせ激減"),
        ("移動中", "音声入力で話すだけ\nナレッジに自動追加"),
        ("夕方（3分）", "AIの提案をレビュー\n承認 or 却下をタップ"),
    ]
    for i, (time, desc) in enumerate(timeline_m):
        x = 0.6 + i * 3.1
        add_textbox(s, x, 2.2, 2.8, 0.5, time,
                    font_size=20, color=IND_COLORS["製造業"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.9, 2.8, 2.0, desc,
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(s, 0.8, 5.3, 11, 0.5,
                     "Q&A例: 「この製品の原価率はどのくらいが妥当？」",
                     font_size=18, color=DARK, bold=True)
    add_paragraph(tf,
                  "→ 「過去3年の類似製品は原価率73%。特殊加工込みで76-78%が妥当です」",
                  font_size=17, color=ACCENT)

    # ── Slide 8: 歯科の活用シーン ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "活用シーン ― 歯科",
                font_size=36, color=IND_COLORS["歯科"], bold=True)
    add_textbox(s, 0.8, 1.3, 6, 0.5,
                "院長の1日（導入後）",
                font_size=24, color=DARK, bold=True)
    timeline_d = [
        ("朝（5分）", "ダッシュボードでAI提案を確認\n「今月のリコール率が\n目標を下回っています」"),
        ("日中", "スタッフの質問はAIが回答\n院長への問い合わせ激減"),
        ("昼休み", "音声入力で話すだけ\nナレッジに自動追加"),
        ("夕方（3分）", "AIの提案をレビュー\n承認 or 却下をタップ"),
    ]
    for i, (time, desc) in enumerate(timeline_d):
        x = 0.6 + i * 3.1
        add_textbox(s, x, 2.2, 2.8, 0.5, time,
                    font_size=20, color=IND_COLORS["歯科"], bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.9, 2.8, 2.0, desc,
                    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(s, 0.8, 5.3, 11, 0.5,
                     "Q&A例: 「この症例の自費の案内基準は？」",
                     font_size=18, color=DARK, bold=True)
    add_paragraph(tf,
                  "→ 「院内ルールでは○○の場合は自費セラミックを第一選択。根拠: 院内マニュアルp.12」",
                  font_size=17, color=ACCENT)

    # ── Slide 9: BPO自動化の価値 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "ナレッジが貯まると → AIが「動ける」ようになる",
                font_size=34, color=NAVY, bold=True)
    add_textbox(s, 0.8, 1.3, 11, 0.6,
                "BPO自動化: 社長が承認するだけで、AIが業務を代行",
                font_size=24, color=DARK, bold=True)

    # Before / After
    tf_before = add_textbox(s, 0.8, 2.2, 5.5, 0.5,
                            "■ 今まで（人がやっていた作業）",
                            font_size=20, color=GRAY, bold=True)
    for item in [
        "毎月の請求データをfreeeに手入力（2時間）",
        "日報をExcelに転記してkintoneに登録（毎日30分）",
        "勤怠データをジョブカンに毎週まとめて入力",
        "受注情報をANDPADと社内Excelに二重入力",
        "レセプトデータの突合チェック（月末3時間）",
    ]:
        add_paragraph(tf_before, f"  ✕  {item}", font_size=15, color=GRAY)

    tf_after = add_textbox(s, 7.0, 2.2, 5.5, 0.5,
                           "■ シャチョツー導入後",
                           font_size=20, color=ACCENT, bold=True)
    for item in [
        "AIが請求データを自動でfreeeに反映",
        "音声日報 → AI構造化 → kintoneに自動登録",
        "勤怠データを自動集計・自動入力",
        "受注情報を1回入力 → 全SaaSに自動連携",
        "レセプトの突合を自動チェック → 差異のみ通知",
    ]:
        add_paragraph(tf_after, f"  ✅  {item}", font_size=15, color=DARK)

    # 数字で見る効果
    tf_num = add_textbox(s, 0.8, 5.6, 11, 0.5,
                         "削減効果（1社あたり）",
                         font_size=20, color=DARK, bold=True)
    add_paragraph(tf_num,
                  "月20〜50時間の手作業を削減  →  金額換算 ¥60,000〜150,000/月  →  人件費の再配分が可能に",
                  font_size=18, color=ACCENT, bold=True)
    add_paragraph(tf_num,
                  "※ 全操作に承認フローあり。AIが勝手に実行することはありません。",
                  font_size=15, color=GRAY)

    # ── Slide 10: Engineer支援の価値 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "さらに先へ — AIが「気づき」「考え」「提案する」",
                font_size=34, color=NAVY, bold=True)
    add_textbox(s, 0.8, 1.3, 11, 0.6,
                "Engineer支援: 暗黙知を発見し、業務改善を自動で提案",
                font_size=24, color=DARK, bold=True)

    eng_features = [
        ("行動推論", ORANGE,
         "社員の行動パターンから\n「暗黙のルール」を自動発見。\n\n例: 「ベテランのAさんは\n雨の日の前日に必ず\n資材を追加発注している」\n→ ルール化して全員に共有"),
        ("要件定義の自動生成", ORANGE,
         "「こういう業務を改善したい」と\n社長が話すだけで、\nAIが要件定義書を作成。\n\n必要な情報が足りなければ\nAIが質問して補完。\n→ そのままエンジニアに渡せる"),
        ("業界ベンチマーク", ORANGE,
         "同業他社（匿名）と比較して\n自社の立ち位置がわかる。\n\n例: 「御社の原価率は業界平均\nより8%高い。主因は○○工程」\n→ 改善ポイントが一目瞭然"),
    ]
    for i, (title, feat_color, desc) in enumerate(eng_features):
        x = 0.6 + i * 4.2
        add_textbox(s, x, 2.2, 3.9, 0.5, title,
                    font_size=22, color=feat_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.9, 3.9, 4.0, desc,
                    font_size=16, color=GRAY)

    add_textbox(s, 0.8, 6.5, 11, 0.5,
                "→ ナレッジが貯まるほど精度が上がる。早く始めた企業ほど恩恵が大きい。",
                font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 11: 料金 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "料金", font_size=36, color=NAVY, bold=True)
    add_textbox(s, 0.8, 1.3, 11, 0.6,
                "全プラン共通: 最初の3ヶ月は月額 ¥30,000",
                font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    make_table(s, [
        ["", "ライト", "スタンダード", "エンタープライズ"],
        ["月額（税抜）", "¥200,000", "¥300,000", "¥500,000"],
        ["エージェント", "1体", "3〜5体", "5体〜無制限"],
        ["ユーザー数", "〜15名", "〜50名", "無制限"],
        ["想定企業", "単一事業/1医院", "複数事業/2-3院", "グループ/5院+"],
    ], 1.5, 2.3, 10, 3.5)
    add_textbox(s, 0.8, 6.2, 11, 0.5,
                "まず3ヶ月、ナレッジを貯めてください。効果を実感してから正式にご検討ください。",
                font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 10: ROI ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "投資対効果（ROI）", font_size=36, color=NAVY, bold=True)
    add_textbox(s, 0.8, 1.3, 6, 0.5,
                "ライトプラン（¥200,000/月）の場合",
                font_size=22, color=DARK, bold=True)
    make_table(s, [
        ["効果項目", "金額換算"],
        ["社長/院長の時間削減（月20-30h）", "¥60,000〜90,000"],
        ["情報検索時間の削減", "¥30,000〜50,000"],
        ["属人化リスクの低減", "プライスレス"],
        ["AIによるリスク回避・改善", "¥50,000〜150,000"],
    ], 0.8, 2.0, 7, 3.5)
    add_textbox(s, 8.5, 2.5, 4, 2,
                "ROI\n1.0〜1.5倍+",
                font_size=40, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 8.5, 4.5, 4, 1.5,
                "合計価値:\n¥140,000〜290,000+/月",
                font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 11: 競合比較 ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "競合比較", font_size=36, color=NAVY, bold=True)
    make_table(s, [
        ["", "シャチョツー", "ChatGPT/Copilot", "経営コンサル"],
        ["御社専用の知識", "○", "×（汎用）", "△（属人的）"],
        ["24時間対応", "○", "○", "×"],
        ["業界テンプレ", "○（3業種対応）", "×", "×"],
        ["能動提案", "○", "×", "△（月1回）"],
        ["SaaS自動実行", "○（将来）", "×", "×"],
        ["月額", "¥200,000〜", "¥3,000/人", "¥300,000〜"],
        ["知識の蓄積", "永続", "リセットされる", "担当者依存"],
    ], 0.8, 1.5, 11.5, 5)

    # ── Slide 12: セキュリティ ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "セキュリティ — 御社のナレッジを守る仕組み",
                font_size=36, color=NAVY, bold=True)
    security = [
        "日本国内のサーバーに保管",
        "銀行レベルの暗号化（AES-256-GCM）",
        "他社データとは完全に壁で分離",
        "AIの学習にデータは一切使われない",
        "当社スタッフも無断でデータにアクセスできない",
        "万一の漏洩時はサイバー保険で補償 + 72時間以内に通知",
        "解約時はデータ完全削除 + 削除証明書を発行",
        "歯科: 要配慮個人情報（病歴等）対応 + 3省2ガイドライン準拠",
    ]
    tf = add_textbox(s, 1.5, 1.6, 10, 0.5, "", font_size=20)
    for item in security:
        add_paragraph(tf, f"✅  {item}", font_size=21, color=DARK, space_before=Pt(12))

    # ── Slide 13: 導入フロー ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "導入フロー", font_size=36, color=NAVY, bold=True)
    steps = [
        ("Step 1", "お申込み\n（5分）", "オンラインで\n即日利用開始"),
        ("Step 2", "テンプレート\n選択（10分）", "建設/製造/歯科\nから選択"),
        ("Step 3", "ナレッジ入力\n（初週）", "話す/入力/\n書類取込"),
        ("Step 4", "社員がQ&A\n（2-3週目）", "社長に聞く代わり\nにAIに聞く"),
        ("Step 5", "能動提案\n（1ヶ月目〜）", "AIがリスク・改善\nを先に教える"),
    ]
    for i, (step, title, desc) in enumerate(steps):
        x = 0.5 + i * 2.5
        add_textbox(s, x, 1.8, 2.2, 0.4, step,
                    font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.5, 2.2, 1.0, title,
                    font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 3.7, 2.2, 1.5, desc,
                    font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(s, 0.8, 5.5, 11, 0.6,
                "→ 3ヶ月後、「これなしでは困る」状態に。",
                font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 14: FAQ ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "よくある質問", font_size=36, color=NAVY, bold=True)

    faqs = [
        ("Q: ChatGPTと何が違うの？",
         "ChatGPT = 汎用的な「物知り」。誰に聞いても同じ答え。\n"
         "シャチョツー = 御社専用の「参謀」。御社のナレッジ・ルール・業界知識に基づいて回答。"),
        ("Q: ITに詳しくないけど使える？",
         "話すだけでOK。スマホでも使えます。専門的な設定は不要です。"),
        ("Q: 他社にうちの情報が見えない？",
         "データベースの仕組みで完全に壁があります。アプリにバグがあっても絶対に見えません。"),
        ("Q: 途中で辞められる？",
         "はい。違約金なし。データはJSON/CSVでエクスポート可能。削除証明書も発行。"),
    ]
    y = 1.4
    for question, answer in faqs:
        tf = add_textbox(s, 0.8, y, 11, 0.4, question,
                         font_size=22, color=DARK, bold=True)
        add_paragraph(tf, answer, font_size=17, color=GRAY)
        y += 1.4

    # ── Slide 15: ロードマップ（BPO・Engineer展開時期） ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(s, 0.8, 0.4, 11, 0.8,
                "今後のロードマップ — 使い続けるほど賢く、できることが増える",
                font_size=32, color=NAVY, bold=True)

    # タイムライン横軸
    phases = [
        ("導入〜3ヶ月", "ナレッジ蓄積期", ACCENT,
         "✅ ナレッジ入力（音声/テキスト/書類）\n"
         "✅ Q&A（社員がAIに質問）\n"
         "✅ 能動提案（リスク・改善を通知）\n"
         "✅ 業界テンプレート（80%完成）\n"
         "✅ ダッシュボード（充足度マップ）"),
        ("4ヶ月目〜", "BPO自動化", GREEN,
         "🔜 SaaS自動操作\n"
         "   kintone / freee / ANDPAD 等\n"
         "🔜 ナレッジに基づくタスク自動実行\n"
         "🔜 承認フロー付き（勝手にはやらない）\n"
         "🔜 ROI自動計測ダッシュボード"),
        ("6ヶ月目〜", "Engineer支援", ORANGE,
         "🔮 行動推論（暗黙知の自動発見）\n"
         "🔮 業務改善の要件定義を自動生成\n"
         "🔮 社内システム改修の提案・実行\n"
         "🔮 業界ベンチマーク（匿名比較）\n"
         "🔮 他業種テンプレート順次追加"),
    ]
    for i, (period, label, phase_color, desc) in enumerate(phases):
        x = 0.5 + i * 4.2
        add_textbox(s, x, 1.5, 3.9, 0.4, period,
                    font_size=22, color=phase_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.0, 3.9, 0.4, label,
                    font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(s, x, 2.6, 3.9, 4.0, desc,
                    font_size=16, color=GRAY)

    add_textbox(s, 0.8, 6.5, 11, 0.5,
                "→ 早く始めた企業ほど、ナレッジが貯まり、AIの精度が高く、自動化の恩恵が大きい。",
                font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Slide 16: CTA ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, NAVY)
    add_textbox(s, 1, 1.5, 11, 1.2,
                "まずは3ヶ月、\n月¥30,000でお試しください。",
                font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    tf = add_textbox(s, 2, 3.5, 9, 0.5, "", font_size=22)
    for item in [
        "✅ 初日から業界テンプレートで即利用可能（建設・製造・歯科）",
        "✅ 3ヶ月で効果を実感",
        "✅ 合わなければそのまま終了（違約金なし）",
    ]:
        add_paragraph(tf, item, font_size=24, color=WHITE, space_before=Pt(12))
    add_textbox(s, 1, 5.8, 11, 0.8,
                "今月中のお申込みでオンボーディング支援（¥100,000）を無料！",
                font_size=24, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

    return prs


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

if __name__ == "__main__":
    prs = generate_unified_deck()
    path = "shachotwo/shachotwo_deck.pptx"
    prs.save(path)
    print(f"✅ {path}  ({len(prs.slides)}枚)")
    print()
    print("構成:")
    print("  Slide 1:  表紙")
    print("  Slide 2:  課題（3業種共通）")
    print("  Slide 3:  シャチョツーとは")
    print("  Slide 4:  3つの価値（業種別例付き）")
    print("  Slide 5:  業界テンプレート一覧（建設・製造・歯科）")
    print("  Slide 6:  活用シーン — 建設業")
    print("  Slide 7:  活用シーン — 製造業")
    print("  Slide 8:  活用シーン — 歯科")
    print("  Slide 9:  BPO自動化の価値")
    print("  Slide 10: Engineer支援の価値")
    print("  Slide 11: 料金")
    print("  Slide 12: ROI")
    print("  Slide 13: 競合比較")
    print("  Slide 14: セキュリティ")
    print("  Slide 15: 導入フロー")
    print("  Slide 16: FAQ")
    print("  Slide 17: ロードマップ（展開時期）")
    print("  Slide 18: CTA")
    print()
    print("Canvaへの取り込み:")
    print("  1. https://www.canva.com →「アップロード」")
    print("  2. pptx をドラッグ&ドロップ → 自動変換")
    print("  3. フォント・色・レイアウトを好みに調整")
