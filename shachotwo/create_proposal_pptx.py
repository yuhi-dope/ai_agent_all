#!/usr/bin/env python3
"""
コンサル会社向け提案書 PPTX生成スクリプト
ターゲット: 責任者クラス（稟議用資料）
デザイン: クリーン・ミニマル・プロフェッショナル
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════════════════════════
# デザインシステム
# ═══════════════════════════════════════════════════════════════
# カラーパレット（プロフェッショナル・信頼感）
PRIMARY = RGBColor(0x0F, 0x1F, 0x3C)       # ダークネイビー
PRIMARY_LIGHT = RGBColor(0x1A, 0x36, 0x6B)  # ミディアムネイビー
ACCENT = RGBColor(0x00, 0x6A, 0xF5)         # ブルー
ACCENT_LIGHT = RGBColor(0xE6, 0xF0, 0xFF)   # 薄青背景
ACCENT_DARK = RGBColor(0x00, 0x4A, 0xBF)    # 濃青
SUCCESS = RGBColor(0x00, 0x96, 0x5B)         # グリーン
SUCCESS_LIGHT = RGBColor(0xE6, 0xF7, 0xF0)  # 薄緑
WARNING = RGBColor(0xF0, 0x8C, 0x00)         # オレンジ
DANGER = RGBColor(0xDC, 0x35, 0x45)          # レッド
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK = RGBColor(0x2D, 0x2D, 0x3F)
GRAY_800 = RGBColor(0x33, 0x37, 0x3D)
GRAY_600 = RGBColor(0x6C, 0x75, 0x7D)
GRAY_400 = RGBColor(0xAD, 0xB5, 0xBD)
GRAY_200 = RGBColor(0xE9, 0xEC, 0xEF)
GRAY_100 = RGBColor(0xF8, 0xF9, 0xFA)
GOLD = RGBColor(0xC5, 0x9D, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ═══════════════════════════════════════════════════════════════
# ヘルパー関数
# ═══════════════════════════════════════════════════════════════

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s

def rounded(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=16, color=DARK, bold=False, align=PP_ALIGN.LEFT, font="Meiryo", anchor=None):
    """テキストボックスを追加"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def add_p(tf, text, sz=14, color=DARK, bold=False, align=PP_ALIGN.LEFT, space=Pt(4), font="Meiryo"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = space
    return p

def header_bar(slide, title, subtitle=None):
    """統一ヘッダーバー"""
    rect(slide, Inches(0), Inches(0), W, Inches(1.1), PRIMARY)
    # 左のアクセントライン
    rect(slide, Inches(0), Inches(1.1), W, Pt(3), ACCENT)
    tb(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
       title, sz=28, color=WHITE, bold=True)
    if subtitle:
        tb(slide, Inches(0.8), Inches(0.65), Inches(10), Inches(0.4),
           subtitle, sz=13, color=GRAY_400)

def footer(slide, num, total=18):
    """統一フッター"""
    rect(slide, Inches(0), Inches(7.15), W, Pt(1), GRAY_200)
    tb(slide, Inches(0.8), Inches(7.15), Inches(5), Inches(0.3),
       "Confidential — シャチョツー（社長2号）ご提案書", sz=9, color=GRAY_400)
    tb(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
       f"{num} / {total}", sz=9, color=GRAY_400, align=PP_ALIGN.RIGHT)

def card(slide, l, t, w, h, fill=WHITE, border=GRAY_200, shadow=True):
    """カード型コンテナ（影付き風）"""
    if shadow:
        # 影（少しずらした薄い四角）
        rect(slide, l + Pt(3), t + Pt(3), w, h, GRAY_200)
    s = rounded(slide, l, t, w, h, fill, border)
    return s

def icon_num(slide, l, t, num, color=ACCENT, size=0.65):
    """数字付き丸アイコン"""
    o = oval(slide, l, t, Inches(size), Inches(size), color)
    otf = o.text_frame
    otf.word_wrap = False
    op = otf.paragraphs[0]
    op.text = str(num)
    op.font.size = Pt(int(size * 24))
    op.font.color.rgb = WHITE
    op.font.bold = True
    op.font.name = "Meiryo"
    op.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════
# Slide 1: 表紙
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, PRIMARY)

# 上部の装飾ライン
rect(s, Inches(0), Inches(0), W, Pt(4), ACCENT)

# ドキュメント種別
tb(s, Inches(1.2), Inches(1.2), Inches(5), Inches(0.5),
   "ご提案書  |  社内検討資料", sz=15, color=GRAY_400)

# メインタイトル
rect(s, Inches(1.2), Inches(2.2), Inches(3.5), Pt(3), ACCENT)
tb(s, Inches(1.2), Inches(2.5), Inches(10), Inches(1.2),
   "AI コンサル強化プラットフォーム\n「シャチョツー」導入のご提案", sz=38, color=WHITE, bold=True)

# サブタイトル
tb(s, Inches(1.2), Inches(4.2), Inches(10), Inches(0.8),
   "— コンサルタント生産性 2倍 × OEMリカーリング収入の実現 —", sz=20, color=RGBColor(0x7E, 0xA8, 0xE0))

# メタ情報
tb(s, Inches(1.2), Inches(5.8), Inches(5), Inches(1.2),
   "提出先:  御社 ご担当者様\n提出元:  シャチョツー株式会社\n日付:    2026年3月", sz=14, color=GRAY_400)

tb(s, Inches(8), Inches(6.3), Inches(4.5), Inches(0.5),
   "CONFIDENTIAL", sz=13, color=GOLD, bold=True, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# Slide 2: 目次（稟議に必須）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "目次")

sections = [
    ("01", "エグゼクティブサマリー", "ご提案の全体像"),
    ("02", "現状の課題認識", "3つの構造的課題"),
    ("03", "ソリューション概要", "シャチョツーで何が変わるか"),
    ("04", "導入効果（定量）", "Before/After・育成効果"),
    ("05", "OEMリカーリングモデル", "コンサル終了後の新収益源"),
    ("06", "料金体系", "プラン・OEMマージン構造"),
    ("07", "投資対効果（ROI）", "費用 vs 効果の定量分析"),
    ("08", "比較検討", "自社開発 vs シャチョツー導入"),
    ("09", "リスクと対策", "想定リスクと軽減策"),
    ("10", "セキュリティ", "データ保護・コンプライアンス"),
    ("11", "導入スケジュール", "5ステップ導入フロー"),
    ("12", "判断基準（稟議用）", "Go/No-Go チェックリスト"),
]

for i, (num, title, desc) in enumerate(sections):
    col = i // 6
    row = i % 6
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 0.9)
    tb(s, x, y, Inches(0.6), Inches(0.5), num, sz=16, color=ACCENT, bold=True)
    tb(s, x + Inches(0.7), y, Inches(3.5), Inches(0.3), title, sz=16, color=DARK, bold=True)
    tb(s, x + Inches(0.7), y + Inches(0.35), Inches(5), Inches(0.3), desc, sz=12, color=GRAY_600)

footer(s, 2)

# ═══════════════════════════════════════════════════════════════
# Slide 3: エグゼクティブサマリー
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "01  エグゼクティブサマリー", "本提案の全体像を1ページで")

# メインメッセージ
msg_box = card(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0), ACCENT_LIGHT, ACCENT)
tf = msg_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "「御社のコンサルタント全員を、御社のトップコンサルタントに近づける」AIプラットフォームです"
p.font.size = Pt(20)
p.font.color.rgb = PRIMARY
p.font.bold = True
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER

# 3つの価値カード
values = [
    ("生産性 2倍", "担当社数\n3-5社 → 8-12社/人", "コンサル稼働率の天井を突破"),
    ("新収益源", "OEMリカーリング\nマージン75%（永続）", "コンサル終了後も月次収入"),
    ("ナレッジ資産化", "退職しても\n知見が残る", "属人化を構造的に解消"),
]
for i, (title, metric, desc) in enumerate(values):
    x = Inches(0.8 + i * 4.0)
    c = card(s, x, Inches(2.9), Inches(3.6), Inches(2.6))
    icon_num(s, x + Inches(0.2), Inches(3.05), i + 1)
    tb(s, x + Inches(1.0), Inches(3.05), Inches(2.5), Inches(0.5),
       title, sz=18, color=PRIMARY, bold=True)
    tb(s, x + Inches(0.3), Inches(3.7), Inches(3.0), Inches(1.0),
       metric, sz=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.3), Inches(4.8), Inches(3.0), Inches(0.5),
       desc, sz=12, color=GRAY_600, align=PP_ALIGN.CENTER)

# 導入条件サマリー
cond_box = card(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.9), GRAY_100)
tf = cond_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "導入条件:  初期費用 ¥0  ｜  お試し ¥30,000/月 × 3ヶ月  ｜  正規 ¥200,000〜/月  ｜  違約金なし  ｜  即日解約可"
p.font.size = Pt(15)
p.font.color.rgb = GRAY_800
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER

footer(s, 3)

# ═══════════════════════════════════════════════════════════════
# Slide 4: 現状の課題認識
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "02  現状の課題認識", "コンサルティング会社が直面する3つの構造的課題")

challenges = [
    ("課題①", "ナレッジが個人に依存", DANGER,
     ["プロジェクト知見がコンサルタント個人のPC・メモに散在",
      "ベテラン退職 ＝ 知見ごと消失（年間退職率 15-25%）",
      "ジュニアの立ち上がりに 6ヶ月〜1年",
      "同じ失敗を繰り返す（学習が組織に蓄積しない）"]),
    ("課題②", "稼働率に構造的天井", WARNING,
     ["コンサルタント1人 ＝ 同時 3〜5社が限界",
      "定型作業（分析・資料作成・質問対応）に時間を消費",
      "売上 ＝ 人数 × 担当社数 → 人数に比例しスケールしない",
      "採用→教育→戦力化のサイクルが遅い"]),
    ("課題③", "コンサル終了後、効果が消える", RGBColor(0x6F, 0x42, 0xC1),
     ["3-6ヶ月のプロジェクト終了後、半年で効果が薄れる",
      "「知識移転」≒ 報告書を渡して終わり",
      "リカーリング収入がない → 毎期ゼロから営業",
      "「いい仕事でしたが、もう大丈夫です」で終了"]),
]

for i, (label, title, accent, items) in enumerate(challenges):
    x = Inches(0.8 + i * 4.1)
    c = card(s, x, Inches(1.5), Inches(3.8), Inches(5.2))
    # アクセントバー上部
    rect(s, x, Inches(1.5), Inches(3.8), Pt(5), accent)
    tb(s, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.4),
       label, sz=13, color=accent, bold=True)
    tb(s, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.5),
       title, sz=18, color=PRIMARY, bold=True)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.3), Inches(3.0 + j * 0.8), Inches(3.3), Inches(0.7),
           f"• {item}", sz=13, color=GRAY_800)

footer(s, 4)

# ═══════════════════════════════════════════════════════════════
# Slide 5: ソリューション概要
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "03  ソリューション概要", "シャチョツーが提供する3つの機能")

solutions = [
    ("ナレッジ資産化", "過去の知見を構造化・検索可能に",
     "過去プロジェクトの知見を即座に検索・引用\n提案書ドラフトの自動生成\n業界ベンチマークをリアルタイム提供\n退職者のナレッジが永続的に組織に残る",
     "→ 課題①を解決"),
    ("コンサル生産性 2倍", "AIが定型作業を代替",
     "クライアントの定型質問にAIが回答\n議事録の自動要約・時系列整理\n類似案件の自動比較・チェック\nジュニアの育成期間を1/3に短縮",
     "→ 課題②を解決"),
    ("OEMリカーリング収入", "コンサル終了後の新収益源",
     "コンサル知見をAIに蓄積 → 終了後も継続利用\n御社ブランドでクライアントに提供（OEM）\nマージン75%が永続的に発生\n開発・保守・セキュリティは当社負担",
     "→ 課題③を解決"),
]

for i, (title, subtitle, desc, link) in enumerate(solutions):
    x = Inches(0.8 + i * 4.1)
    c = card(s, x, Inches(1.5), Inches(3.8), Inches(5.0))
    icon_num(s, x + Inches(0.25), Inches(1.7), i + 1, ACCENT)
    tb(s, x + Inches(1.1), Inches(1.75), Inches(2.6), Inches(0.4),
       title, sz=18, color=PRIMARY, bold=True)
    tb(s, x + Inches(0.3), Inches(2.3), Inches(3.3), Inches(0.4),
       subtitle, sz=13, color=GRAY_600)
    rect(s, x + Inches(0.3), Inches(2.7), Inches(3.2), Pt(1), GRAY_200)
    tb(s, x + Inches(0.3), Inches(2.9), Inches(3.3), Inches(2.5),
       desc, sz=13, color=GRAY_800)
    # 課題リンク
    tb(s, x + Inches(0.3), Inches(5.7), Inches(3.3), Inches(0.4),
       link, sz=13, color=ACCENT, bold=True)

footer(s, 5)

# ═══════════════════════════════════════════════════════════════
# Slide 6: 導入効果（定量）Before/After
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "04  導入効果（定量）", "主要KPIの Before / After 比較")

# KPI比較テーブル
metrics = [
    ("コンサルタント1人の担当社数", "3-5社", "8-12社", "2〜3倍"),
    ("コンサルタント1人の月間売上", "¥150-250万", "¥400-600万", "2倍以上"),
    ("提案書作成工数", "40時間/月", "15時間/月", "▲63%"),
    ("ジュニア育成期間", "6ヶ月", "2ヶ月", "1/3"),
    ("コンサル終了後のリカーリング", "¥0", "¥15万/月/社", "新規収益"),
    ("退職時のナレッジ損失", "全損失", "ゼロ（AIに蓄積）", "100%防止"),
]

col_w = [Inches(4.0), Inches(2.5), Inches(2.5), Inches(2.0)]
headers_text = ["指標", "Before", "After", "改善幅"]
ht = Inches(0.55)
y0 = Inches(1.5)

# ヘッダー行
for j, (h, cw) in enumerate(zip(headers_text, col_w)):
    x = Inches(0.8) + sum(col_w[:j])
    rect(s, x, y0, cw, ht, PRIMARY)
    tb(s, x, y0 + Inches(0.08), cw, Inches(0.4),
       h, sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# データ行
for i, (metric, before, after, delta) in enumerate(metrics):
    y = y0 + ht * (i + 1)
    bg = WHITE if i % 2 == 0 else GRAY_100
    vals = [metric, before, after, delta]
    colors = [GRAY_800, DANGER, SUCCESS, ACCENT]
    bolds = [False, False, True, True]
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER]
    for j, (v, cw) in enumerate(zip(vals, col_w)):
        x = Inches(0.8) + sum(col_w[:j])
        rect(s, x, y, cw, ht, bg)
        lpad = Inches(0.3) if j == 0 else Inches(0)
        tb(s, x + lpad, y + Inches(0.08), cw - lpad, Inches(0.4),
           v, sz=14, color=colors[j], bold=bolds[j], align=aligns[j])

# 注記
tb(s, Inches(0.8), Inches(5.2), Inches(10), Inches(0.5),
   "※ 担当社数UPは「コンサルタントの時間が空く→新規受注可能」として試算。売上は業界平均単価で算出。",
   sz=11, color=GRAY_600)

footer(s, 6)

# ═══════════════════════════════════════════════════════════════
# Slide 7: ジュニア育成効果
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "04  導入効果（育成加速）", "ジュニアコンサルタントの育成期間を1/3に短縮")

# 従来タイムライン
tb(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
   "従来の育成プロセス", sz=18, color=DANGER, bold=True)

bar_left = Inches(0.8)
bar_w_old = Inches(10)
for i in range(6):
    x = bar_left + Inches(i * 1.667)
    clr = RGBColor(0xFF, 0xE0, 0xE0) if i < 5 else RGBColor(0xFF, 0xCC, 0xCC)
    rect(s, x, Inches(2.1), Inches(1.667), Inches(0.6), clr, GRAY_200)
    tb(s, x, Inches(2.15), Inches(1.667), Inches(0.5),
       f"{i+1}ヶ月目", sz=11, color=GRAY_600, align=PP_ALIGN.CENTER)
tb(s, Inches(11), Inches(2.15), Inches(1.5), Inches(0.5),
   "独り立ち", sz=14, color=DANGER, bold=True)

tb(s, Inches(0.8), Inches(2.9), Inches(10), Inches(0.4),
   "ベテランOJT → 資料読み込み → 案件見学 → 少しずつ担当 → 6ヶ月でようやく1人対応可能",
   sz=13, color=GRAY_600)

# 新タイムライン
tb(s, Inches(0.8), Inches(3.7), Inches(5), Inches(0.4),
   "シャチョツー導入後", sz=18, color=SUCCESS, bold=True)

for i in range(2):
    x = bar_left + Inches(i * 1.667)
    rect(s, x, Inches(4.3), Inches(1.667), Inches(0.6), SUCCESS_LIGHT, SUCCESS)
    tb(s, x, Inches(4.35), Inches(1.667), Inches(0.5),
       f"{i+1}ヶ月目", sz=11, color=GRAY_600, align=PP_ALIGN.CENTER)
tb(s, Inches(4.2), Inches(4.35), Inches(1.5), Inches(0.5),
   "独り立ち", sz=14, color=SUCCESS, bold=True)

# ジュニアがAIに聞ける例
tb(s, Inches(0.8), Inches(5.3), Inches(5), Inches(0.4),
   "ジュニアがAIに質問できる例:", sz=15, color=PRIMARY, bold=True)

examples = [
    "「この業界の標準的な利益率は？」→ AIが業界データで即回答",
    "「前任者はこのクライアントにどう対応してた？」→ 過去の対応履歴から回答",
    "「この提案で抜けてる観点は？」→ 過去の類似提案と比較してチェック",
]
for j, ex in enumerate(examples):
    tb(s, Inches(1.0), Inches(5.8 + j * 0.4), Inches(11), Inches(0.4),
       ex, sz=13, color=GRAY_800)

footer(s, 7)

# ═══════════════════════════════════════════════════════════════
# Slide 8: OEMリカーリングモデル
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "05  OEMリカーリングモデル", "コンサル終了後の新収益源 — Salesforceモデル準拠")

# 従来 vs 新モデル
for i, (label, color, body) in enumerate([
    ("従来モデル", DANGER,
     "プロジェクト: ¥300万 × 6ヶ月 = ¥1,800万\n終了後: ¥0（毎期ゼロから営業）"),
    ("シャチョツー導入後", SUCCESS,
     "プロジェクト: ¥1,800万 + 継続リカーリング\nOEMマージン: ¥15万/月/社（永続）"),
]):
    x = Inches(0.8 + i * 6.4)
    c = card(s, x, Inches(1.5), Inches(5.8), Inches(1.8))
    rect(s, x, Inches(1.5), Inches(5.8), Pt(4), color)
    tb(s, x + Inches(0.3), Inches(1.7), Inches(5), Inches(0.4),
       label, sz=16, color=color, bold=True)
    tb(s, x + Inches(0.3), Inches(2.2), Inches(5.2), Inches(0.9),
       body, sz=14, color=GRAY_800)

# OEMフロー図
flow_box = card(s, Inches(0.8), Inches(3.7), Inches(11.7), Inches(1.5), ACCENT_LIGHT, ACCENT)
tf = flow_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "OEM収益フロー"
p.font.size = Pt(14)
p.font.color.rgb = ACCENT
p.font.bold = True
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER
add_p(tf, "", sz=4)
add_p(tf, "[クライアント] ─ 月額¥200,000〜 →  [御社]  ─ OEMライセンス料（25%） →  [シャチョツー]",
      sz=18, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
add_p(tf, "御社マージン: 請求額の 75%（永続）  ｜  開発・保守・インフラ・セキュリティは全て当社が負担",
      sz=13, color=GRAY_600, align=PP_ALIGN.CENTER)

# ボリュームディスカウント
tb(s, Inches(0.8), Inches(5.5), Inches(5), Inches(0.4),
   "ボリュームディスカウント（Salesforce準拠）", sz=16, color=PRIMARY, bold=True)

vol_headers = ["OEM経由の月間売上", "当社取り分", "御社マージン"]
vol_rows = [
    ("〜¥500万/月", "25%", "75%"),
    ("¥500万〜2,000万", "20%", "80%"),
    ("¥2,000万〜", "15%", "85%"),
]
vol_cw = [Inches(3.0), Inches(1.8), Inches(1.8)]
vy0 = Inches(5.9)
for j, (h, cw) in enumerate(zip(vol_headers, vol_cw)):
    x = Inches(0.8) + sum(vol_cw[:j])
    rect(s, x, vy0, cw, Inches(0.4), PRIMARY)
    tb(s, x, vy0 + Pt(3), cw, Inches(0.35), h, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(vol_rows):
    y = vy0 + Inches(0.4 + i * 0.35)
    bg = WHITE if i % 2 == 0 else GRAY_100
    for j, (v, cw) in enumerate(zip(row, vol_cw)):
        x = Inches(0.8) + sum(vol_cw[:j])
        rect(s, x, y, cw, Inches(0.35), bg)
        clr = SUCCESS if j == 2 else GRAY_800
        tb(s, x, y + Pt(2), cw, Inches(0.3), v, sz=12, color=clr, bold=(j==2), align=PP_ALIGN.CENTER)

# 右側に10社シミュレーション
tb(s, Inches(7.5), Inches(5.5), Inches(5), Inches(0.4),
   "10社提供時のシミュレーション", sz=16, color=PRIMARY, bold=True)

sim = [
    ("月間OEMマージン", "¥1,500,000/月"),
    ("年間リカーリング", "¥18,000,000"),
    ("LTV（平均30ヶ月・10社）", "¥45,000,000"),
]
for i, (label, val) in enumerate(sim):
    y = Inches(6.0 + i * 0.35)
    tb(s, Inches(7.5), y, Inches(3.2), Inches(0.35), label, sz=13, color=GRAY_800)
    tb(s, Inches(10.7), y, Inches(2), Inches(0.35), val, sz=14, color=SUCCESS, bold=True, align=PP_ALIGN.RIGHT)

footer(s, 8)

# ═══════════════════════════════════════════════════════════════
# Slide 9: コンサルタントの1日
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "05  具体的な活用イメージ", "御社コンサルタントの1日")

timeline = [
    ("朝 5分", "ダッシュボード確認", "担当クライアントのリスクアラートを確認。AIが原因仮説を3つ提示。そのままクライアントに報告。", ACCENT),
    ("午前", "ミーティング準備", "「A社の過去3回のMTG要点をまとめて」→ 即座に要約。同業種のベストプラクティスも提示。", SUCCESS),
    ("午後", "提案書作成", "過去の成功提案書 + A社固有データ → ドラフト自動生成。作成工数 1/3に。", WARNING),
    ("夕方", "新人サポート", "新人は「シャチョツーに聞いて」で自己解決。ベテランの時間を消費しない。", RGBColor(0x6F, 0x42, 0xC1)),
]

for i, (time, title, desc, accent) in enumerate(timeline):
    y = Inches(1.5 + i * 1.4)
    # タイムバッジ
    badge = rounded(s, Inches(0.8), y, Inches(1.4), Inches(0.5), accent)
    btf = badge.text_frame
    bp = btf.paragraphs[0]
    bp.text = time
    bp.font.size = Pt(13)
    bp.font.color.rgb = WHITE
    bp.font.bold = True
    bp.font.name = "Meiryo"
    bp.alignment = PP_ALIGN.CENTER
    # タイトル
    tb(s, Inches(2.5), y - Inches(0.05), Inches(4), Inches(0.45),
       title, sz=17, color=PRIMARY, bold=True)
    # 説明
    tb(s, Inches(2.5), y + Inches(0.45), Inches(10), Inches(0.8),
       desc, sz=14, color=GRAY_800)

footer(s, 9)

# ═══════════════════════════════════════════════════════════════
# Slide 10: 料金体系
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "06  料金体系", "社内利用プラン")

plans = [
    ("ライト", "¥200,000", "/月", "1体", "〜15名", "専門特化コンサル", False),
    ("スタンダード", "¥300,000", "/月", "3-5体", "〜50名", "中堅コンサル", True),
    ("エンタープライズ", "¥500,000", "/月", "5体〜", "無制限", "総合コンサル", False),
]

for i, (name, price, unit, agents, users, target, recommended) in enumerate(plans):
    x = Inches(0.8 + i * 4.1)
    border = ACCENT if recommended else GRAY_200
    c = card(s, x, Inches(1.5), Inches(3.8), Inches(4.5), WHITE, border)
    if recommended:
        rect(s, x, Inches(1.5), Inches(3.8), Pt(4), ACCENT)
        badge = rounded(s, x + Inches(1.0), Inches(1.35), Inches(1.8), Inches(0.35), ACCENT)
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.text = "推奨"
        bp.font.size = Pt(11)
        bp.font.color.rgb = WHITE
        bp.font.bold = True
        bp.font.name = "Meiryo"
        bp.alignment = PP_ALIGN.CENTER

    tb(s, x + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.5),
       name, sz=20, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.3), Inches(2.6), Inches(3.2), Inches(0.6),
       price, sz=30, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.3), Inches(3.1), Inches(3.2), Inches(0.3),
       unit, sz=14, color=GRAY_600, align=PP_ALIGN.CENTER)

    details = [
        f"エージェント数:  {agents}",
        f"利用ユーザー:    {users}",
        f"想定規模:       {target}",
    ]
    for j, d in enumerate(details):
        tb(s, x + Inches(0.5), Inches(3.7 + j * 0.5), Inches(3), Inches(0.4),
           d, sz=13, color=GRAY_800)

# お試し条件
trial = card(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6), RGBColor(0xFF, 0xF9, 0xE6), WARNING)
tf = trial.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "全プラン共通:  初期3ヶ月は ¥30,000/月（お試し）  →  合わなければ終了OK  ｜  違約金なし"
p.font.size = Pt(15)
p.font.color.rgb = GRAY_800
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER

footer(s, 10)

# ═══════════════════════════════════════════════════════════════
# Slide 11: ROI分析（稟議の核心）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "07  投資対効果（ROI）", "コンサルタント10名の会社 — スタンダードプラン適用時")

# 投資サイド
tb(s, Inches(0.8), Inches(1.4), Inches(2), Inches(0.4),
   "年間投資額", sz=16, color=DANGER, bold=True)

inv_box = card(s, Inches(0.8), Inches(1.9), Inches(4.5), Inches(1.8), RGBColor(0xFF, 0xF5, 0xF5), DANGER)
tf = inv_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "月額 ¥300,000 × 12ヶ月"
p.font.size = Pt(16)
p.font.color.rgb = GRAY_800
p.font.name = "Meiryo"
add_p(tf, "年間投資額: ¥3,600,000", sz=24, color=DANGER, bold=True, space=Pt(12))
add_p(tf, "※ お試し期間3ヶ月は¥30,000/月", sz=11, color=GRAY_600, space=Pt(8))

# 効果サイド
tb(s, Inches(6.0), Inches(1.4), Inches(2), Inches(0.4),
   "年間効果額", sz=16, color=SUCCESS, bold=True)

eff_box = card(s, Inches(6.0), Inches(1.9), Inches(6.5), Inches(1.8), SUCCESS_LIGHT, SUCCESS)
tf = eff_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ""
p.font.size = Pt(2)
p.font.name = "Meiryo"

effects = [
    ("担当社数UP（3社→5社/人×10人）", "+¥48,000,000/年"),
    ("提案書作成時間削減（40h→15h/月）", "+¥6,000,000/年"),
    ("新人育成期間短縮（6ヶ月→2ヶ月）", "+¥800,000/人"),
]
for label, val in effects:
    add_p(tf, f"  {label}    {val}", sz=13, color=GRAY_800, space=Pt(6))

add_p(tf, "年間効果額: ¥54,000,000+", sz=24, color=SUCCESS, bold=True, space=Pt(10))

# ROIハイライト
roi_box = card(s, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.5), ACCENT, ACCENT)
tf = roi_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "ROI:  約15倍"
p.font.size = Pt(40)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER
add_p(tf, "年間投資 ¥360万  →  年間効果 ¥5,400万+  ｜  投資回収期間: 約1ヶ月", sz=16, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# 注記
tb(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
   "※ 上記にOEMリカーリング収入（¥1,800万/年 ※10社時）は含まず。OEM併用時のROIはさらに拡大。\n※ 担当社数UPは新規受注可能数として試算。実際の受注は市場環境・営業力に依存。",
   sz=11, color=GRAY_600)

footer(s, 11)

# ═══════════════════════════════════════════════════════════════
# Slide 12: 比較検討（自社開発 vs 導入）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "08  比較検討", "自社開発 vs シャチョツー導入")

comp_h = ["比較項目", "自社開発", "シャチョツー導入", "備考"]
comp_r = [
    ("初期費用", "¥5,000万〜", "¥0", "開発人件費+インフラ"),
    ("開発期間", "12-18ヶ月", "即日利用開始", "POC含む"),
    ("月額ランニング", "¥200万〜", "¥30万/月", "人件費+インフラ+保守"),
    ("エンジニア採用", "2-3名必要", "不要", "AI+インフラ+セキュリティ"),
    ("セキュリティ保守", "自社責任", "当社が24/365で負担", "RLS監査含む"),
    ("業種テンプレート", "ゼロから構築", "50業種対応済み", "構築に数年相当"),
    ("ベンチマーク精度", "自社30社分のみ", "500社+のデータ", "統計的有意性"),
    ("機能アップデート", "自社で都度開発", "自動反映（無料）", "月次リリース"),
]

col_w2 = [Inches(2.8), Inches(2.8), Inches(3.0), Inches(2.4)]
ht2 = Inches(0.5)
y0_2 = Inches(1.5)

for j, (h, cw) in enumerate(zip(comp_h, col_w2)):
    x = Inches(0.8) + sum(col_w2[:j])
    rect(s, x, y0_2, cw, ht2, PRIMARY)
    tb(s, x, y0_2 + Inches(0.07), cw, Inches(0.4), h, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row_data in enumerate(comp_r):
    y = y0_2 + ht2 * (i + 1)
    bg = WHITE if i % 2 == 0 else GRAY_100
    for j, (v, cw) in enumerate(zip(row_data, col_w2)):
        x = Inches(0.8) + sum(col_w2[:j])
        rect(s, x, y, cw, ht2, bg)
        if j == 2:
            clr, bld = SUCCESS, True
        elif j == 1:
            clr, bld = GRAY_800, False
        else:
            clr, bld = GRAY_800, (j == 0)
        lpad = Inches(0.2) if j == 0 else Inches(0)
        tb(s, x + lpad, y + Inches(0.07), cw - lpad, Inches(0.4), v, sz=12, color=clr, bold=bld, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# 3年TCO比較
tb(s, Inches(0.8), Inches(5.8), Inches(5), Inches(0.4),
   "3年間TCO（Total Cost of Ownership）比較", sz=16, color=PRIMARY, bold=True)

tco_items = [
    ("自社開発", "¥5,000万 + ¥200万×36ヶ月 = ¥12,200万", DANGER),
    ("シャチョツー", "¥30万×36ヶ月 = ¥1,080万", SUCCESS),
    ("差額", "▲¥11,120万（91%削減）", ACCENT),
]
for i, (label, val, clr) in enumerate(tco_items):
    y = Inches(6.3 + i * 0.35)
    tb(s, Inches(1.0), y, Inches(2), Inches(0.35), label, sz=14, color=clr, bold=True)
    tb(s, Inches(3.2), y, Inches(8), Inches(0.35), val, sz=14, color=GRAY_800)

footer(s, 12)

# ═══════════════════════════════════════════════════════════════
# Slide 13: リスクと対策（稟議に必須）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "09  リスクと対策", "想定されるリスクと軽減策")

risks = [
    ("コンサルタントが使いこなせない", "低",
     "UIは「質問するだけ」の直感操作。導入初期にオンボーディング支援（対面/オンライン）を提供。まず1チーム試用→全社展開の段階導入。"),
    ("AIの回答精度が不十分", "中",
     "全回答に引用元を表示（検証可能）。フィードバック機能で精度向上。3ヶ月の¥30,000お試し期間で品質を確認後に正規移行。"),
    ("クライアントデータの漏洩", "低",
     "DB層で完全テナント分離（RLS）。AES-256暗号化。AI学習にデータ不使用（契約保証）。AI駆動セキュリティで常時監査。"),
    ("ベンダーロックイン", "低",
     "データはJSON/CSVで一括エクスポート可能。違約金なし・即日解約可。SLA契約明記。"),
    ("シャチョツーの事業継続性", "中",
     "サイバー保険加入済み。解約後30日のエクスポート猶予+削除証明書。SLA契約で品質保証。データのポータビリティを確保。"),
]

risk_h = ["リスク項目", "発生可能性", "軽減策"]
risk_cw = [Inches(3.0), Inches(1.2), Inches(7.5)]
ry0 = Inches(1.5)

for j, (h, cw) in enumerate(zip(risk_h, risk_cw)):
    x = Inches(0.5) + sum(risk_cw[:j])
    rect(s, x, ry0, cw, Inches(0.45), PRIMARY)
    tb(s, x, ry0 + Inches(0.05), cw, Inches(0.35), h, sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, (risk, prob, mitigation) in enumerate(risks):
    y = ry0 + Inches(0.45 + i * 1.0)
    bg = WHITE if i % 2 == 0 else GRAY_100
    prob_color = SUCCESS if prob == "低" else WARNING
    for j, (v, cw) in enumerate(zip([risk, prob, mitigation], risk_cw)):
        x = Inches(0.5) + sum(risk_cw[:j])
        rect(s, x, y, cw, Inches(1.0), bg)
    tb(s, Inches(0.7), y + Inches(0.15), Inches(2.6), Inches(0.7), risk, sz=13, color=GRAY_800, bold=True)
    # 可能性バッジ
    prob_badge = rounded(s, Inches(3.7), y + Inches(0.25), Inches(0.8), Inches(0.35), prob_color)
    ptf = prob_badge.text_frame
    pp = ptf.paragraphs[0]
    pp.text = prob
    pp.font.size = Pt(11)
    pp.font.color.rgb = WHITE
    pp.font.bold = True
    pp.font.name = "Meiryo"
    pp.alignment = PP_ALIGN.CENTER
    tb(s, Inches(4.9), y + Inches(0.1), Inches(7.3), Inches(0.8), mitigation, sz=12, color=GRAY_800)

footer(s, 13)

# ═══════════════════════════════════════════════════════════════
# Slide 14: セキュリティ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "10  セキュリティ", "クライアントのデータを守る多層防御")

sec_categories = [
    ("データ保護", [
        ("日本国内サーバー", "データは国内に保管"),
        ("AES-256暗号化", "保存時暗号化（銀行・軍レベル）"),
        ("TLS 1.3", "通信時暗号化"),
    ]),
    ("アクセス制御", [
        ("完全テナント分離", "DB層RLSで構造的に分離"),
        ("ロールベースアクセス", "CEO→Employee 5段階権限"),
        ("JWT認証", "15分有効期限+リフレッシュ"),
    ]),
    ("AI安全性", [
        ("学習にデータ不使用", "契約上保証（DPA締結）"),
        ("PII自動検出+マスク", "正規表現+LLMハイブリッド"),
        ("全操作監査ログ", "いつ誰が何をしたか記録"),
    ]),
    ("運用保守", [
        ("AI駆動セキュリティ", "RLS監査・脆弱性スキャン自動"),
        ("解約時完全削除", "削除証明書を発行"),
        ("サイバー保険加入済み", "万が一の事故に備え"),
    ]),
]

for i, (cat, items) in enumerate(sec_categories):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.4)
    y = Inches(1.5 + row * 2.8)
    c = card(s, x, y, Inches(5.8), Inches(2.5))
    tb(s, x + Inches(0.3), y + Inches(0.15), Inches(5), Inches(0.4),
       cat, sz=16, color=PRIMARY, bold=True)
    rect(s, x + Inches(0.3), y + Inches(0.55), Inches(5.2), Pt(1), GRAY_200)
    for j, (item, desc) in enumerate(items):
        iy = y + Inches(0.7 + j * 0.55)
        tb(s, x + Inches(0.3), iy, Inches(0.3), Inches(0.4), "✓", sz=14, color=SUCCESS, bold=True)
        tb(s, x + Inches(0.6), iy, Inches(2.2), Inches(0.4), item, sz=13, color=GRAY_800, bold=True)
        tb(s, x + Inches(2.9), iy, Inches(2.7), Inches(0.4), desc, sz=12, color=GRAY_600)

footer(s, 14)

# ═══════════════════════════════════════════════════════════════
# Slide 15: 導入スケジュール
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "11  導入スケジュール", "5ステップの段階導入 — リスクを最小化")

steps = [
    ("STEP 1", "お試し導入", "3ヶ月", "¥30,000/月", [
        "1チームで試用開始",
        "過去の知見を10件投入",
        "AI回答品質を検証",
    ]),
    ("STEP 2", "社内浸透", "1-2ヶ月", "正規料金", [
        "全コンサルタントに展開",
        "知見を網羅的に投入",
        "テンプレート活用開始",
    ]),
    ("STEP 3", "効果検証", "1ヶ月", "—", [
        "生産性変化を測定",
        "ROIダッシュボード確認",
        "OEM提供の可否判断",
    ]),
    ("STEP 4", "OEM準備", "1ヶ月", "—", [
        "御社ブランド設計",
        "クライアント提案資料",
        "OEM契約締結",
    ]),
    ("STEP 5", "OEM開始", "本格展開", "—", [
        "既存客への提案開始",
        "リカーリング計上開始",
        "新規セット提案",
    ]),
]

colors_step = [ACCENT, RGBColor(0x00, 0x56, 0xD2), RGBColor(0x00, 0x42, 0xAA), RGBColor(0x00, 0x30, 0x85), PRIMARY_LIGHT]

for i, (step, title, period, cost, items) in enumerate(steps):
    x = Inches(0.3 + i * 2.55)
    # ステップカード
    c = card(s, x, Inches(1.5), Inches(2.4), Inches(5.0), WHITE, colors_step[i])
    rect(s, x, Inches(1.5), Inches(2.4), Pt(4), colors_step[i])

    # ステップ番号
    tb(s, x + Inches(0.2), Inches(1.7), Inches(2.0), Inches(0.4),
       step, sz=12, color=colors_step[i], bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.2), Inches(2.1), Inches(2.0), Inches(0.4),
       title, sz=16, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.3), Inches(2.55), Inches(1.8), Pt(1), GRAY_200)
    tb(s, x + Inches(0.2), Inches(2.7), Inches(2.0), Inches(0.3),
       period, sz=13, color=colors_step[i], bold=True, align=PP_ALIGN.CENTER)
    if cost != "—":
        tb(s, x + Inches(0.2), Inches(3.0), Inches(2.0), Inches(0.3),
           cost, sz=12, color=GRAY_600, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.3), Inches(3.5 + j * 0.5), Inches(1.9), Inches(0.45),
           f"• {item}", sz=11, color=GRAY_800)

    # 矢印（最後以外）
    if i < 4:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.35), Inches(3.5), Inches(0.25), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY_400
        arrow.line.fill.background()

footer(s, 15)

# ═══════════════════════════════════════════════════════════════
# Slide 16: Salesforceモデル参考
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "[参考]  Salesforceモデルとの構造比較", "Salesforce × Accenture の成功パターンを再現")

# 左: Salesforce
c1 = card(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5))
tb(s, Inches(1.0), Inches(1.7), Inches(5), Inches(0.4),
   "Salesforce × Accenture", sz=18, color=PRIMARY, bold=True)
rect(s, Inches(1.0), Inches(2.1), Inches(5.0), Pt(1), GRAY_200)
sf_steps = [
    "Step 1: CRMツールとしてコンサル会社に提供",
    "Step 2: コンサル会社が「導入支援」を新サービス化",
    "Step 3: クライアント業務がSalesforce前提に",
    "Step 4: Salesforceなしでは提案不可に",
]
for j, step in enumerate(sf_steps):
    tb(s, Inches(1.2), Inches(2.3 + j * 0.5), Inches(5), Inches(0.4),
       step, sz=13, color=GRAY_800)
tb(s, Inches(1.2), Inches(4.4), Inches(5), Inches(0.4),
   "→ 年商¥5兆の40%+がパートナー経由", sz=13, color=ACCENT, bold=True)

# 右: シャチョツー
c2 = card(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.5))
tb(s, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4),
   "シャチョツー × 御社", sz=18, color=PRIMARY, bold=True)
rect(s, Inches(7.2), Inches(2.1), Inches(5.0), Pt(1), GRAY_200)
st_steps = [
    "Step 1: コンサル強化ツールとして御社に提供",
    "Step 2: 御社が「ナレッジDXパッケージ」を新サービス化",
    "Step 3: クライアントのナレッジがAIに蓄積",
    "Step 4: 御社の提案にシャチョツーが標準組込",
]
for j, step in enumerate(st_steps):
    tb(s, Inches(7.4), Inches(2.3 + j * 0.5), Inches(5), Inches(0.4),
       step, sz=13, color=GRAY_800)
tb(s, Inches(7.4), Inches(4.4), Inches(5), Inches(0.4),
   "→ 売上UP + リカーリング + 顧客満足度UP", sz=13, color=SUCCESS, bold=True)

# マージン比較
tb(s, Inches(0.8), Inches(5.3), Inches(5), Inches(0.4),
   "マージン構造比較", sz=16, color=PRIMARY, bold=True)

margin_h = ["モデル", "プラットフォーム取り分", "パートナー取り分"]
margin_r = [
    ("Salesforce OEM", "25%", "75%"),
    ("Salesforce ISV", "15%", "85%"),
    ("シャチョツー OEM", "25%", "75%"),
    ("シャチョツー（高ボリューム時）", "15%", "85%"),
]
mcw = [Inches(4), Inches(3), Inches(3)]
my0 = Inches(5.7)
for j, (h, cw) in enumerate(zip(margin_h, mcw)):
    x = Inches(0.8) + sum(mcw[:j])
    rect(s, x, my0, cw, Inches(0.35), PRIMARY)
    tb(s, x, my0 + Pt(2), cw, Inches(0.3), h, sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, row in enumerate(margin_r):
    y = my0 + Inches(0.35 + i * 0.3)
    bg = WHITE if i % 2 == 0 else GRAY_100
    for j, (v, cw) in enumerate(zip(row, mcw)):
        x = Inches(0.8) + sum(mcw[:j])
        rect(s, x, y, cw, Inches(0.3), bg)
        clr = SUCCESS if j == 2 else GRAY_800
        tb(s, x, y + Pt(1), cw, Inches(0.25), v, sz=11, color=clr, bold=(j==2), align=PP_ALIGN.CENTER)

footer(s, 16)

# ═══════════════════════════════════════════════════════════════
# Slide 17: 判断基準（稟議用チェックリスト）
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
header_bar(s, "12  判断基準（稟議用）", "Go / No-Go 判断チェックリスト")

# Go条件
tb(s, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
   "Go条件（以下を満たす場合、導入を推奨）", sz=16, color=SUCCESS, bold=True)

go_items = [
    "コンサルタントの担当社数に天井を感じている",
    "ベテラン退職時のナレッジ流出が経営課題",
    "コンサル終了後のリカーリング収入を確保したい",
    "年間 ¥360万（月¥30万）の投資が可能",
    "まず1チーム3ヶ月のトライアルに社内合意が取れる",
    "AIツールの社内導入に経営層のサポートがある",
]

for i, item in enumerate(go_items):
    y = Inches(1.9 + i * 0.5)
    tb(s, Inches(1.0), y, Inches(0.4), Inches(0.4), "□", sz=16, color=SUCCESS, bold=True)
    tb(s, Inches(1.5), y, Inches(5.5), Inches(0.4), item, sz=14, color=GRAY_800)

# No-Go条件
tb(s, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
   "No-Go条件（以下に該当する場合は見送り）", sz=16, color=DANGER, bold=True)

nogo_items = [
    "コンサルタント全員がITツールに強い抵抗がある",
    "ナレッジを組織資産にする意思がない",
    "月¥30,000のお試し費用も負担できない",
    "3ヶ月以内に効果を出す意思がない",
]

for i, item in enumerate(nogo_items):
    y = Inches(1.9 + i * 0.5)
    tb(s, Inches(7.2), y, Inches(0.4), Inches(0.4), "□", sz=16, color=DANGER, bold=True)
    tb(s, Inches(7.7), y, Inches(5), Inches(0.4), item, sz=14, color=GRAY_800)

# 稟議サマリーボックス
summary_box = card(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), ACCENT_LIGHT, ACCENT)
tf = summary_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "稟議サマリー"
p.font.size = Pt(16)
p.font.color.rgb = PRIMARY
p.font.bold = True
p.font.name = "Meiryo"

summary_lines = [
    ("件名:", "AI コンサル強化プラットフォーム「シャチョツー」導入"),
    ("投資額:", "初期¥0 ｜ お試し¥30,000/月×3ヶ月 ｜ 正規¥300,000/月（スタンダード）"),
    ("期待効果:", "コンサル生産性2倍（年間+¥5,400万） + OEMリカーリング（年間+¥1,800万）"),
    ("リスク:", "3ヶ月お試し＋違約金なしにより、投資リスクは最大¥90,000に限定"),
    ("推奨:", "まず1チーム（¥30,000/月×3ヶ月）でPoCを実施し、効果を検証後に全社展開"),
]
for label, val in summary_lines:
    add_p(tf, f"  {label}  {val}", sz=13, color=GRAY_800, space=Pt(6))

footer(s, 17)

# ═══════════════════════════════════════════════════════════════
# Slide 18: Next Step (CTA)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, PRIMARY)
rect(s, Inches(0), Inches(0), W, Pt(4), ACCENT)

tb(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
   "Next Step", sz=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

cta = card(s, Inches(2.5), Inches(2.8), Inches(8.3), Inches(3.3), WHITE)
tf = cta.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "まずは30分のデモ + ヒアリングから"
p.font.size = Pt(26)
p.font.color.rgb = PRIMARY
p.font.bold = True
p.font.name = "Meiryo"
p.alignment = PP_ALIGN.CENTER

cta_steps = [
    "① 御社のコンサルスタイル・課題をヒアリング（15分）",
    "② 実際のデモで動作をお見せ（10分）",
    "③ 御社に合った導入プランをご提案（5分）",
]
for step in cta_steps:
    add_p(tf, step, sz=18, color=GRAY_800, align=PP_ALIGN.CENTER, space=Pt(10))

add_p(tf, "", sz=8)
add_p(tf, "お試し: ¥30,000/月 × 3ヶ月  ｜  初期費用: ¥0  ｜  違約金: なし", sz=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, space=Pt(12))

tb(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
   "お問い合わせ:  [ご担当者連絡先]", sz=16, color=RGBColor(0x7E, 0xA8, 0xE0), align=PP_ALIGN.CENTER)

tb(s, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
   "CONFIDENTIAL  —  シャチョツー株式会社", sz=11, color=GRAY_600, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════════════
output_path = "/Users/sugimotoyuuhi/code/ai_agent/shachotwo/コンサル会社向け提案書.pptx"
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
